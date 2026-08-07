import streamlit as st
import pandas as pd
from pandas.api.types import is_numeric_dtype
import numpy as np
import matplotlib.pyplot as plt
from io import BytesIO
import zipfile
import textwrap
import joblib
import time
from auth_ui import render_auth_page
from auth import (
    authenticate_user,
    create_user,
    initialize_auth_database,
    validate_signup,
)
initialize_auth_database()

if "authenticated" not in st.session_state:
    st.session_state["authenticated"] = False

if "current_user" not in st.session_state:
    st.session_state["current_user"] = None

if "auth_mode" not in st.session_state:
    st.session_state["auth_mode"] = "Login"

if not st.session_state["authenticated"]:
    render_auth_page()
    st.stop()    

user = st.session_state["current_user"]

st.markdown(
    f"""
    <h2>👋 Welcome back, {user['full_name']}!</h2>
    <p style="color:gray;">
        Ready to build your next machine learning model?
    </p>
    """,
    unsafe_allow_html=True,
)
try:
    import shap
except Exception:
    shap = None

from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pptx import Presentation

from sklearn.model_selection import train_test_split, cross_val_score, StratifiedKFold, RandomizedSearchCV
from sklearn.ensemble import RandomForestClassifier, GradientBoostingClassifier, ExtraTreesClassifier, AdaBoostClassifier
from sklearn.tree import DecisionTreeClassifier
from sklearn.svm import SVC
from sklearn.neighbors import KNeighborsClassifier
from sklearn.naive_bayes import GaussianNB
from sklearn.pipeline import make_pipeline
from sklearn.preprocessing import StandardScaler
from sklearn.linear_model import LogisticRegression
from sklearn.metrics import accuracy_score, confusion_matrix, classification_report, roc_curve, auc, f1_score
from sklearn.inspection import permutation_importance
from sklearn.base import clone

# Optional AutoML models. LabMind will use them if installed, and skip them safely if not.
try:
    from xgboost import XGBClassifier
except Exception:
    XGBClassifier = None

try:
    from lightgbm import LGBMClassifier
except Exception:
    LGBMClassifier = None

try:
    from catboost import CatBoostClassifier
except Exception:
    CatBoostClassifier = None


st.set_page_config(page_title="LabMind.ai", page_icon="🧠", layout="wide")


def create_pdf_report(report_text):
    buffer = BytesIO()
    pdf = canvas.Canvas(buffer, pagesize=letter)
    width, height = letter

    pdf.setFont("Helvetica-Bold", 20)
    pdf.drawString(40, height - 50, "LabMind.ai Analysis Report")

    pdf.setFont("Helvetica", 10)
    y = height - 90

    for line in report_text.split("\n"):
        if y < 50:
            pdf.showPage()
            pdf.setFont("Helvetica", 10)
            y = height - 50

        pdf.drawString(40, y, line[:100])
        y -= 15

    pdf.save()
    buffer.seek(0)
    return buffer


def create_ppt_report(report_text):
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "LabMind.ai Analysis Report"
    slide.placeholders[1].text = "AI-powered dataset analysis and machine learning report"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    slide.placeholders[1].text = report_text[:1200]

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Recommendations"
    slide.placeholders[1].text = (
        "• Review missing-value-heavy columns\n"
        "• Remove duplicate rows\n"
        "• Avoid ID columns as predictors\n"
        "• Use the best model as a baseline\n"
        "• Validate on external datasets before deployment"
    )

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


def auto_clean_dataset(df):
    cleaned_df = df.copy()
    cleaned_df = cleaned_df.drop_duplicates()

    for col in cleaned_df.columns:
        if cleaned_df[col].dtype == "object":
            cleaned_df[col] = cleaned_df[col].fillna("Unknown")
        else:
            cleaned_df[col] = pd.to_numeric(cleaned_df[col], errors="coerce")
            median_value = cleaned_df[col].median()
            cleaned_df[col] = cleaned_df[col].fillna(
                0 if pd.isna(median_value) else median_value
            )

    return cleaned_df


def clean_uploaded_dataframe(df):
    cleaned = df.copy()

    # Drop columns like Unnamed: 32 and completely empty columns.
    cleaned = cleaned.loc[:, ~cleaned.columns.astype(str).str.contains(r"^Unnamed", case=False, regex=True)]
    cleaned = cleaned.dropna(axis=1, how="all")

    return cleaned


def reset_model_state_if_new_file(uploaded_file):
    file_signature = f"{uploaded_file.name}_{getattr(uploaded_file, 'size', 'unknown')}"

    if st.session_state.get("uploaded_file_signature") != file_signature:
        keys_to_clear = [
            "best_model",
            "best_model_name",
            "best_accuracy",
            "target_column",
            "feature_columns",
            "category_maps",
            "numeric_defaults",
            "target_label_map",
            "report",
            "executive_summary",
            "cleaned_df",
            "trained_models",
            "leaderboard",
            "classification_report_df",
            "confusion_matrix",
            "roc_auc",
            "explain_df",
            "explain_method",
            "prediction_history",
            "cleaning_steps",
            "X_test",
            "y_test",
            "X_train",
            "y_train",
            "model_cards",
            "tuning_results",
        ]

        for key in keys_to_clear:
            st.session_state.pop(key, None)

        st.session_state["uploaded_file_signature"] = file_signature


def detect_default_target(df):
    preferred_targets = [
        "diagnosis",
        "target",
        "outcome",
        "label",
        "class",
        "survived",
        "churn",
        "heart_disease",
        "death_event",
    ]

    lower_columns = {str(col).lower(): col for col in df.columns}

    for target_name in preferred_targets:
        if target_name in lower_columns:
            return list(df.columns).index(lower_columns[target_name])

    # If no obvious target exists, choose the last column instead of an ID column.
    for idx in range(len(df.columns) - 1, -1, -1):
        col = str(df.columns[idx]).lower()
        if col not in ["id", "patientid", "patient_id", "passengerid"]:
            return idx

    return 0


def prepare_ml_data(df, target_column):
    data = df.copy()

    # Drop obvious ID / text identifier columns, but NEVER drop the selected target.
    drop_cols = ["PassengerId", "Name", "Ticket", "Cabin", "ID", "Id", "id"]
    drop_cols = [c for c in drop_cols if c in data.columns and c != target_column]
    data = data.drop(columns=drop_cols, errors="ignore")

    category_maps = {}
    numeric_defaults = {}
    target_label_map = None

    # -----------------------------
    # IMPORTANT FIX:
    # Do NOT decide target handling using dtype == "object" only.
    # Some CSVs load string columns as pandas StringDtype / category, not plain object.
    # If we send diagnosis M/B through pd.to_numeric(errors="coerce"), it becomes all NaN,
    # then all 0, causing the "only one class" bug.
    # -----------------------------
    target_raw = data[target_column]

    if is_numeric_dtype(target_raw):
        target_numeric = pd.to_numeric(target_raw, errors="coerce")

        # If numeric conversion accidentally destroys most values, treat target as categorical.
        if target_numeric.notna().sum() < max(2, int(0.5 * len(target_raw))):
            target_series = target_raw.fillna("Unknown").astype(str).str.strip().astype("category")
            target_categories = list(target_series.cat.categories)
            target_label_map = {int(i): str(label) for i, label in enumerate(target_categories)}
            data[target_column] = target_series.cat.codes.astype(int)
        else:
            target_median = target_numeric.median()
            target_filled = target_numeric.fillna(0 if pd.isna(target_median) else target_median)

            # IMPORTANT MULTICLASS FIX:
            # XGBoost requires class labels to be encoded as 0..n_classes-1.
            # Wine Quality targets often look like 3,4,5,6,7,8, which makes XGBoost fail.
            # For classification-like numeric targets, encode labels safely and keep a map
            # so the UI still shows the original labels.
            unique_values = sorted(pd.Series(target_filled.unique()).dropna().tolist())
            if 2 <= len(unique_values) <= 20:
                target_categories = unique_values
                reverse_map = {value: int(i) for i, value in enumerate(target_categories)}
                data[target_column] = target_filled.map(reverse_map).astype(int)
                target_label_map = {int(i): value for i, value in enumerate(target_categories)}
            else:
                data[target_column] = target_filled
    else:
        target_series = target_raw.fillna("Unknown").astype(str).str.strip().astype("category")
        target_categories = list(target_series.cat.categories)
        target_label_map = {int(i): str(label) for i, label in enumerate(target_categories)}
        data[target_column] = target_series.cat.codes.astype(int)

    # Process feature columns only after target is safely handled.
    for col in data.columns:
        if col == target_column:
            continue

        if is_numeric_dtype(data[col]):
            data[col] = pd.to_numeric(data[col], errors="coerce")
            median_value = data[col].median()
            fill_value = 0 if pd.isna(median_value) else median_value
            numeric_defaults[col] = float(fill_value)
            data[col] = data[col].fillna(fill_value)
        else:
            data[col] = data[col].fillna("Unknown").astype(str).str.strip()
            cats = list(data[col].astype("category").cat.categories)
            category_maps[col] = cats
            data[col] = data[col].astype("category").cat.codes.astype(int)

    data = data.replace([float("inf"), float("-inf")], 0)
    data = data.fillna(0)

    X = data.drop(columns=[target_column])
    y = data[target_column]

    return data, X, y, category_maps, numeric_defaults, target_label_map

def get_model_importance(model, feature_names):
    """Return feature importances or coefficients for plain estimators and sklearn pipelines."""
    estimator = model

    if hasattr(model, "steps"):
        estimator = model.steps[-1][1]

    if hasattr(estimator, "feature_importances_"):
        return estimator.feature_importances_

    if hasattr(estimator, "coef_"):
        coef = estimator.coef_
        if len(coef.shape) == 1:
            return abs(coef)
        return abs(coef).mean(axis=0)

    return [0] * len(feature_names)




def decode_target_value(value, target_label_map):
    """Decode encoded target values back to original labels when available."""
    if not target_label_map:
        return value
    if value in target_label_map:
        return target_label_map[value]
    try:
        int_value = int(value)
        if int_value in target_label_map:
            return target_label_map[int_value]
    except Exception:
        pass
    return value


def make_readable_classification_report(y_test, predictions, target_label_map):
    """Create a classification report with human-readable class labels."""
    report = classification_report(y_test, predictions, output_dict=True, zero_division=0)
    report_df = pd.DataFrame(report).transpose()

    readable_index = []
    for idx in report_df.index:
        try:
            if str(idx).replace('.', '', 1).isdigit():
                decoded = decode_target_value(float(idx), target_label_map)
                if decoded == float(idx):
                    decoded = decode_target_value(int(float(idx)), target_label_map)
                readable_index.append(str(decoded))
            else:
                readable_index.append(str(idx))
        except Exception:
            readable_index.append(str(idx))

    report_df.index = readable_index
    return report_df

def apply_dark_chart_theme(fig, ax=None):
    """Apply LabMind dark styling to a Matplotlib figure."""
    fig.patch.set_facecolor("#111827")

    axes = fig.get_axes() if ax is None else [ax]

    for axis in axes:
        axis.set_facecolor("#0f172a")

        axis.title.set_color("white")
        axis.xaxis.label.set_color("white")
        axis.yaxis.label.set_color("white")

        axis.tick_params(axis="x", colors="white")
        axis.tick_params(axis="y", colors="white")

        for spine in axis.spines.values():
            spine.set_color("#475569")

        axis.grid(True, alpha=0.15)

        legend = axis.get_legend()
        if legend is not None:
            legend.get_frame().set_facecolor("#111827")
            legend.get_frame().set_edgecolor("#475569")

            for text in legend.get_texts():
                text.set_color("white")

    fig.tight_layout()
    return fig

def render_probability_bars(probabilities, classes, target_label_map):
    """Render prediction probabilities as clean horizontal bars."""
    st.write("### Prediction Probability")
    prob_rows = []
    for cls, prob in zip(classes, probabilities):
        prob_rows.append({
            "Class": str(decode_target_value(cls, target_label_map)),
            "Probability": float(prob)
        })

    prob_df = pd.DataFrame(prob_rows).sort_values("Probability", ascending=False)
    st.dataframe(
        prob_df.assign(Probability=lambda d: d["Probability"].map(lambda x: f"{x * 100:.2f}%")),
        use_container_width=True
    )

    fig_prob, ax_prob = plt.subplots(figsize=(8, max(2.8, 0.55 * len(prob_df))))
    ax_prob.barh(prob_df["Class"], prob_df["Probability"] * 100)
    ax_prob.set_xlabel("Probability (%)")
    ax_prob.set_title("Prediction Confidence by Class")
    ax_prob.set_xlim(0, 100)
    ax_prob.invert_yaxis()
    apply_dark_chart_theme(fig_prob)
    st.pyplot(fig_prob)


def get_inner_estimator(model):
    """Return final estimator from a sklearn pipeline or the model itself."""
    if hasattr(model, "steps"):
        return model.steps[-1][1]
    return model


def compute_shap_importance(model, X_background, X_sample, feature_names):
    """Try SHAP global importance safely. Returns None when unsupported or slow/failing."""
    if shap is None:
        return None

    estimator = get_inner_estimator(model)

    # SHAP is safest and fastest for tree-style models.
    if not hasattr(estimator, "feature_importances_"):
        return None

    try:
        explainer = shap.TreeExplainer(estimator)
        values = explainer.shap_values(X_sample)

        if isinstance(values, list):
            arr = np.mean([np.abs(v) for v in values], axis=0)
        else:
            arr = np.abs(values)
            if arr.ndim == 3:
                arr = arr.mean(axis=2)

        mean_abs = arr.mean(axis=0)
        shap_df = pd.DataFrame({
            "Feature": list(feature_names),
            "Mean |SHAP value|": mean_abs
        }).sort_values("Mean |SHAP value|", ascending=False)
        return shap_df
    except Exception:
        return None


def compute_universal_explainability(model, X_train, X_test, y_test, feature_names):
    """Explain any winning model with the best available method.
    Priority: SHAP for supported tree models, native importances/coefficients, then permutation importance.
    """
    estimator = get_inner_estimator(model)
    sample = X_test.sample(min(80, len(X_test)), random_state=42) if len(X_test) > 80 else X_test.copy()

    # 1) SHAP for supported tree-style models.
    shap_df = compute_shap_importance(model, X_train, sample, feature_names)
    if shap_df is not None and not shap_df.empty:
        shap_df = shap_df.rename(columns={"Mean |SHAP value|": "Importance"})
        shap_df["Explanation"] = shap_df["Feature"].apply(
            lambda x: f"{x} had strong SHAP impact on model predictions."
        )
        return shap_df.sort_values("Importance", ascending=False), "SHAP Explainability", "SHAP estimates how much each feature pushes predictions up or down."

    # 2) Native feature importance / logistic coefficients.
    native_vals = None
    native_method = None
    if hasattr(estimator, "feature_importances_"):
        native_vals = estimator.feature_importances_
        native_method = "Native Feature Importance"
    elif hasattr(estimator, "coef_"):
        coef = estimator.coef_
        if len(coef.shape) == 1:
            native_vals = np.abs(coef)
        else:
            native_vals = np.abs(coef).mean(axis=0)
        native_method = "Coefficient Importance"

    if native_vals is not None and len(native_vals) == len(feature_names):
        native_df = pd.DataFrame({"Feature": list(feature_names), "Importance": native_vals})
        native_df = native_df.sort_values("Importance", ascending=False)
        native_df["Explanation"] = native_df["Feature"].apply(
            lambda x: f"{x} influenced the model according to {native_method.lower()}."
        )
        return native_df, native_method, "This method uses the model's built-in importance or learned coefficients."

    # 3) Model-agnostic fallback: permutation importance.
    try:
        perm = permutation_importance(
            model,
            X_test,
            y_test,
            n_repeats=5,
            random_state=42,
            scoring="accuracy",
        )
        perm_df = pd.DataFrame({
            "Feature": list(feature_names),
            "Importance": perm.importances_mean,
        }).sort_values("Importance", ascending=False)
        perm_df["Explanation"] = perm_df["Feature"].apply(
            lambda x: f"When {x} was shuffled, model performance changed."
        )
        return perm_df, "Permutation Importance", "Permutation importance works for any model by testing performance drop after shuffling each feature."
    except Exception as e:
        fallback = pd.DataFrame({"Feature": list(feature_names), "Importance": [0] * len(feature_names)})
        fallback["Explanation"] = "No reliable explainability score was available for this model."
        return fallback, "Fallback Explainability", f"LabMind could not compute model-specific explanations: {e}"

def build_classification_models(y):
    """Build a stronger AutoML model list with safe optional models."""
    models = {
        "Random Forest": RandomForestClassifier(
            n_estimators=300, random_state=42, class_weight="balanced"
        ),
        "Extra Trees": ExtraTreesClassifier(
            n_estimators=300, random_state=42, class_weight="balanced"
        ),
        "Gradient Boosting": GradientBoostingClassifier(random_state=42),
        "AdaBoost": AdaBoostClassifier(random_state=42),
        "Decision Tree": DecisionTreeClassifier(random_state=42, class_weight="balanced"),
        "Logistic Regression": make_pipeline(
            StandardScaler(),
            LogisticRegression(max_iter=4000, class_weight="balanced")
        ),
        "SVM": make_pipeline(
            StandardScaler(),
            SVC(probability=True, class_weight="balanced", random_state=42)
        ),
        "KNN": make_pipeline(
            StandardScaler(),
            KNeighborsClassifier(n_neighbors=5)
        ),
        "Naive Bayes": GaussianNB(),
    }

    if XGBClassifier is not None:
        models["XGBoost"] = XGBClassifier(
            n_estimators=250,
            learning_rate=0.05,
            max_depth=4,
            random_state=42,
            eval_metric="logloss",
            verbosity=0,
        )

    if LGBMClassifier is not None:
        models["LightGBM"] = LGBMClassifier(
            n_estimators=250,
            learning_rate=0.05,
            random_state=42,
            verbose=-1,
        )

    if CatBoostClassifier is not None:
        models["CatBoost"] = CatBoostClassifier(
            iterations=250,
            learning_rate=0.05,
            depth=5,
            random_seed=42,
            verbose=False,
        )

    return models


def generate_recommendations(df, health_score, missing_percentage, duplicate_rows):
    recs = []

    if missing_percentage > 20:
        recs.append("High missing data detected. Clean missing values before relying on model results.")
    elif missing_percentage > 0:
        recs.append("Some missing data detected. Median/mode filling is reasonable for a baseline.")
    else:
        recs.append("No missing values detected.")

    if duplicate_rows > 0:
        recs.append(f"Remove {duplicate_rows} duplicate rows before final modeling.")
    else:
        recs.append("No duplicate rows detected.")

    missing_cols = df.isnull().sum().sort_values(ascending=False)
    missing_cols = missing_cols[missing_cols > 0]

    for col, count in missing_cols.items():
        pct = (count / len(df)) * 100
        if pct > 40:
            recs.append(f"{col} has {pct:.2f}% missing values. Consider dropping this column.")
        else:
            recs.append(f"{col} has {pct:.2f}% missing values. Consider imputing it.")

    if health_score >= 85:
        recs.append("Dataset quality is strong enough for baseline modeling.")
    elif health_score >= 65:
        recs.append("Dataset quality is moderate. Clean before serious modeling.")
    else:
        recs.append("Dataset quality is weak. Cleaning should be prioritized.")

    return recs


def generate_executive_summary(
    df,
    health_score,
    missing_percentage,
    duplicate_rows,
    best_model_name=None,
    best_accuracy=None,
    target_column=None,
    top_feature=None
):
    summary = f"""
LabMind analyzed a dataset with {df.shape[0]} rows and {df.shape[1]} columns.

The dataset health score is {health_score}/100. Missing values represent {missing_percentage:.2f}% of the dataset, and {duplicate_rows} duplicate rows were detected.
"""

    if best_model_name is not None:
        summary += f"""
AutoML selected {best_model_name} as the best-performing model for predicting {target_column}. The model achieved an accuracy of {best_accuracy * 100:.2f}%.
"""

    if top_feature is not None:
        summary += f"""
The strongest feature influencing the model was {top_feature}. This means LabMind found this column most useful for prediction.
"""

    summary += """
Recommended next steps:
- Clean high-missing columns.
- Avoid using ID columns as predictors.
- Validate the model on a second dataset before deployment.
- Export the model only after confirming performance.
"""

    return summary



def build_feature_ai_summary(explain_df, method_name, best_model_name, target_column):
    """Create a polished, non-repetitive summary of the most important model drivers."""
    top_features = explain_df.head(5).copy()
    names = top_features["Feature"].astype(str).tolist()

    if len(names) == 0:
        return "LabMind could not identify clear feature drivers for this model."

    if len(names) == 1:
        feature_phrase = names[0]
    elif len(names) == 2:
        feature_phrase = f"{names[0]} and {names[1]}"
    else:
        feature_phrase = ", ".join(names[:-1]) + f", and {names[-1]}"

    return (
        f"LabMind used {method_name.lower()} to explain the {best_model_name} model. "
        f"The strongest drivers for predicting {target_column} were {feature_phrase}. "
        "These features contributed the most to the model's decision-making, so they should be reviewed first when interpreting predictions or validating the model."
    )


def prepare_top_importance_table(explain_df, top_n=10):
    """Return a clean top-N importance table with rank and percentage contribution."""
    clean = explain_df.copy().sort_values("Importance", ascending=False).head(top_n)
    total = clean["Importance"].sum()
    if total and total > 0:
        clean["Importance %"] = clean["Importance"].apply(lambda x: f"{(x / total) * 100:.2f}%")
    else:
        clean["Importance %"] = "N/A"
    clean["Importance"] = clean["Importance"].apply(lambda x: round(float(x), 5))
    medals = ["🥇", "🥈", "🥉"] + [f"#{i}" for i in range(4, top_n + 1)]
    clean.insert(0, "Rank", medals[:len(clean)])
    return clean[["Rank", "Feature", "Importance", "Importance %"]]


def render_confusion_summary(cm):
    """Display quick confusion-matrix insights for binary and multiclass settings."""
    correct = int(np.trace(cm))
    total = int(cm.sum())
    incorrect = total - correct
    acc = (correct / total * 100) if total else 0

    c1, c2, c3 = st.columns(3)
    c1.metric("Correct Predictions", correct)
    c2.metric("Incorrect Predictions", incorrect)
    c3.metric("Matrix Accuracy", f"{acc:.2f}%")

    if cm.shape == (2, 2):
        tn, fp, fn, tp = cm.ravel()
        b1, b2, b3, b4 = st.columns(4)
        b1.metric("True Negative", int(tn))
        b2.metric("False Positive", int(fp))
        b3.metric("False Negative", int(fn))
        b4.metric("True Positive", int(tp))


def render_auc_summary(roc_auc):
    """Add a plain-English interpretation for AUC."""
    if roc_auc >= 0.95:
        label = "Excellent discrimination"
        text = "The model separates the two classes extremely well."
    elif roc_auc >= 0.85:
        label = "Strong discrimination"
        text = "The model separates the two classes well."
    elif roc_auc >= 0.70:
        label = "Moderate discrimination"
        text = "The model has useful signal, but could be improved."
    else:
        label = "Weak discrimination"
        text = "The model may not separate classes reliably yet."

    st.markdown(f"""
    <div class="model-highlight">
        📈 <b>AUC Score: {roc_auc:.3f}</b><br>
        <b>{label}</b>. {text}
    </div>
    """, unsafe_allow_html=True)



def create_professional_pdf_report(report_text, leaderboard_df=None, explain_df=None, prediction_history=None):
    """Create a richer executive PDF report with key sections."""
    buffer = BytesIO()
    pdf = canvas.Canvas(buffer, pagesize=letter)
    width, height = letter

    def new_page(title=None):
        pdf.showPage()
        if title:
            pdf.setFont("Helvetica-Bold", 16)
            pdf.drawString(40, height - 45, title)
            return height - 75
        return height - 50

    pdf.setFont("Helvetica-Bold", 24)
    pdf.drawString(40, height - 60, "LabMind.ai Executive Report")
    pdf.setFont("Helvetica", 12)
    pdf.drawString(40, height - 88, "AI-powered dataset analysis, AutoML, explainability, and deployment summary")
    y = height - 125

    pdf.setFont("Helvetica-Bold", 14)
    pdf.drawString(40, y, "Executive Summary")
    y -= 24
    pdf.setFont("Helvetica", 9)
    for line in report_text.split("\n"):
        if y < 50:
            y = new_page("Executive Summary Continued")
            pdf.setFont("Helvetica", 9)
        pdf.drawString(40, y, line[:115])
        y -= 13

    if leaderboard_df is not None and len(leaderboard_df) > 0:
        y = new_page("Model Leaderboard")
        pdf.setFont("Helvetica", 10)
        for _, row in leaderboard_df.head(12).iterrows():
            if y < 50:
                y = new_page("Model Leaderboard Continued")
                pdf.setFont("Helvetica", 10)
            pdf.drawString(50, y, f"{row.get('Model', '')}: Accuracy {row.get('Accuracy', '')}, CV Mean {row.get('CV Mean', 'N/A')}")
            y -= 16

    if explain_df is not None and len(explain_df) > 0:
        y = new_page("Top Model Drivers")
        pdf.setFont("Helvetica", 10)
        for _, row in explain_df.head(10).iterrows():
            if y < 50:
                y = new_page("Top Model Drivers Continued")
                pdf.setFont("Helvetica", 10)
            pdf.drawString(50, y, f"{row.get('Feature', '')}: {float(row.get('Importance', 0)):.5f}")
            y -= 16

    if prediction_history is not None and len(prediction_history) > 0:
        y = new_page("Prediction History")
        pdf.setFont("Helvetica", 10)
        for row in prediction_history[-10:]:
            if y < 50:
                y = new_page("Prediction History Continued")
                pdf.setFont("Helvetica", 10)
            pdf.drawString(50, y, f"{row.get('Model', '')} predicted {row.get('Target', '')} = {row.get('Prediction', '')} ({row.get('Confidence', 'N/A')})")
            y -= 16

    pdf.save()
    buffer.seek(0)
    return buffer


def get_class_balance_text(y, target_label_map=None):
    counts = pd.Series(y).value_counts().sort_index()
    rows = []
    for cls, count in counts.items():
        rows.append(f"- {decode_target_value(cls, target_label_map)}: {int(count)}")
    return "\n".join(rows)


def answer_dataset_question(q, df, numeric_cols, text_cols, health_score, total_missing, missing_percentage, duplicate_rows, recommendations):
    q = q.lower()
    if "row" in q:
        return f"This dataset has {df.shape[0]} rows."
    if "missing" in q:
        ms = df.isnull().sum()
        ms = ms[ms > 0].sort_values(ascending=False)
        if len(ms) == 0:
            return "There are no missing values in the current cleaned dataset."
        return "Columns with missing values:\n\n" + "\n".join([f"- {c}: {v}" for c, v in ms.items()])
    if "duplicate" in q:
        return f"This dataset has {duplicate_rows} duplicate rows."
    if "health" in q or "quality" in q:
        return f"The dataset health score is {health_score}/100. Missing data is {missing_percentage:.2f}% and duplicate rows are {duplicate_rows}."
    if "numeric" in q:
        return f"The numeric columns are:\n\n{numeric_cols}"
    if "text" in q or "categor" in q:
        return f"The text/category columns are:\n\n{text_cols}"
    if "remove" in q or "drop" in q or "clean" in q:
        return "Recommended cleaning actions:\n\n" + "\n".join([f"- {r}" for r in recommendations])
    if "best model" in q or "deploy" in q:
        model = st.session_state.get("best_model_name")
        acc = st.session_state.get("best_accuracy")

        if model is None or acc is None:
            return (
                "No model has been trained yet. "
                "Train models in the AutoML tab first."
            )

        return (
            f"The current best model is {model} with "
            f"{acc * 100:.2f}% accuracy. "
            "Use it as a baseline, then validate it with "
            "cross-validation and external data before deployment."
        )

    if "best accuracy" in q or "highest accuracy" in q:
        model = st.session_state.get("best_model_name")
        acc = st.session_state.get("best_accuracy")

        if acc is None:
            return (
                "No accuracy result is available yet. "
                "Train models in the AutoML tab first."
            )

        model_text = f" by {model}" if model else ""

        return (
            f"The best accuracy is {acc * 100:.2f}%"
            f"{model_text}."
        )

    if "accuracy" in q:
        acc = st.session_state.get("best_accuracy")
        model = st.session_state.get("best_model_name")

        if acc is None:
            return (
                "No accuracy result is available yet. "
                "Train models first."
            )

        return (
            f"The current best validation accuracy is "
            f"{acc * 100:.2f}% from {model}."
        )

    if "f1" in q or "precision" in q or "recall" in q:
        if "classification_report_df" in st.session_state:
            return (
                "Current model metrics:\n\n"
                + st.session_state[
                    "classification_report_df"
                ].round(4).to_string()
            )

        return (
            "No classification report is available yet. "
            "Train a model first."
        )
    if "feature" in q or "important" in q or "explain" in q:
        if "explain_df" in st.session_state:
            top = st.session_state["explain_df"].head(8)
            return "Top model drivers:\n\n" + "\n".join([f"- {r.Feature}: {r.Importance:.5f}" for r in top.itertuples()])
        return "No explainability results are available yet. Train a model first."
    if "imbalanc" in q or "balance" in q:
        if "target_column" in st.session_state:
            return "Class balance is available after training. Check the AutoML results and classification report for support counts."
        return "Choose a target column and train a model first to inspect class balance."
    if "summary" in q or "summarize" in q:
        return f"""This dataset has {df.shape[0]} rows and {df.shape[1]} columns.

Health score: {health_score}/100.
Missing values: {total_missing}.
Missing percentage: {missing_percentage:.2f}%.
Duplicate rows: {duplicate_rows}.
Numeric columns: {len(numeric_cols)}.
Text/category columns: {len(text_cols)}.
"""
    if "why" in q and ("win" in q or "best" in q or "outperform" in q):
        model = st.session_state.get("best_model_name")
        acc = st.session_state.get("best_accuracy")
        exp = st.session_state.get("model_explainability_summary", "Feature importance is available after model training.")
        if model and acc is not None:
            return f"{model} is currently leading because it achieved the highest validation accuracy ({acc * 100:.2f}%). {exp}"
        return "Train models first, then LabMind can explain why the winning model performed best."
    if "correl" in q:
        if len(numeric_cols) < 2:
            return "There are not enough numeric columns to compute correlations."
        corr = df[numeric_cols].corr().abs().unstack().sort_values(ascending=False)
        corr = corr[corr < 1].head(8)
        return "Strongest numeric correlations:\n\n" + "\n".join([f"- {a} and {b}: {v:.3f}" for (a, b), v in corr.items()])
    if "doctor" in q or "clinician" in q or "medical" in q:
        return f"Clinical-style summary: the dataset contains {df.shape[0]} records and {df.shape[1]} columns. The current target is {st.session_state.get('target_column', 'not selected')}. The best model is {st.session_state.get('best_model_name', 'not trained yet')} with accuracy {st.session_state.get('best_accuracy', 0) * 100:.2f}% if trained. This should be treated as a research/demo model only and externally validated before any real clinical use."
    if "remove" in q and "what happens" in q:
        return "To test feature removal safely, create a copy of the dataset without that column and retrain. If accuracy and cross-validation stay stable, the feature may not be critical. If performance drops, the feature is important."

    return "Try asking about rows, missing values, class balance, cleaning, best model, accuracy, feature importance, correlations, why a model won, deployment, or summary."


def create_model_zip(trained_models):
    buffer = BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as zf:
        for name, bundle in trained_models.items():
            model_buffer = BytesIO()
            joblib.dump(bundle["model"], model_buffer)
            model_buffer.seek(0)
            safe_name = name.lower().replace(" ", "_").replace("/", "_")
            zf.writestr(f"{safe_name}.pkl", model_buffer.read())
    buffer.seek(0)
    return buffer


def create_cleaning_steps(df, removed_columns, duplicate_rows, missing_percentage):
    steps = []
    if removed_columns:
        steps.append(f"Removed empty/system columns: {removed_columns}")
    else:
        steps.append("No empty/system columns needed removal")
    steps.append(f"Detected {duplicate_rows} duplicate rows")
    steps.append(f"Detected {missing_percentage:.2f}% missing data")
    steps.append("Auto Clean can fill numeric missing values with median")
    steps.append("Auto Clean can fill text missing values with Unknown")
    steps.append("Categorical features are encoded safely during model training")
    steps.append("ID-like columns are avoided as predictors when possible")
    return steps


def compare_two_datasets(df_a, df_b, name_a="Current dataset", name_b="Comparison dataset"):
    rows = []
    for name, data in [(name_a, df_a), (name_b, df_b)]:
        total_cells = data.shape[0] * data.shape[1]
        missing = int(data.isnull().sum().sum())
        rows.append({
            "Dataset": name,
            "Rows": data.shape[0],
            "Columns": data.shape[1],
            "Missing Values": missing,
            "Missing %": round((missing / total_cells) * 100, 2) if total_cells else 0,
            "Duplicate Rows": int(data.duplicated().sum()),
            "Numeric Columns": len(data.select_dtypes(include=["number"]).columns),
            "Text Columns": len(data.select_dtypes(include=["object"]).columns),
        })
    return pd.DataFrame(rows)


def tune_top_models(results, trained_models, X_train, y_train, X_test, y_test, top_k=3):
    """Lightweight hyperparameter optimization for the best baseline models."""
    if len(results) == 0:
        return results, trained_models, []

    base_results = pd.DataFrame(results).sort_values("Accuracy", ascending=False).head(top_k)
    tuning_log = []

    param_spaces = {
        "Random Forest": {
            "n_estimators": [150, 250, 400],
            "max_depth": [None, 4, 6, 10, 14],
            "min_samples_split": [2, 4, 8],
            "min_samples_leaf": [1, 2, 4],
        },
        "Extra Trees": {
            "n_estimators": [150, 250, 400],
            "max_depth": [None, 4, 6, 10, 14],
            "min_samples_split": [2, 4, 8],
            "min_samples_leaf": [1, 2, 4],
        },
        "Gradient Boosting": {
            "n_estimators": [80, 120, 180, 250],
            "learning_rate": [0.02, 0.05, 0.08, 0.1],
            "max_depth": [2, 3, 4],
        },
        "AdaBoost": {
            "n_estimators": [50, 100, 150, 250, 350],
            "learning_rate": [0.03, 0.05, 0.1, 0.5, 1.0],
        },
        "Decision Tree": {
            "max_depth": [None, 3, 4, 6, 10, 14],
            "min_samples_split": [2, 4, 8, 12],
            "min_samples_leaf": [1, 2, 4, 8],
        },
        "Logistic Regression": {
            "logisticregression__C": [0.01, 0.05, 0.1, 0.5, 1, 2, 5, 10],
        },
        "SVM": {
            "svc__C": [0.1, 0.5, 1, 2, 5, 10],
            "svc__gamma": ["scale", "auto", 0.01, 0.05, 0.1],
            "svc__kernel": ["rbf", "linear"],
        },
        "KNN": {
            "kneighborsclassifier__n_neighbors": [3, 5, 7, 9, 11],
            "kneighborsclassifier__weights": ["uniform", "distance"],
        },
    }

    results_by_model = {r["Model"]: dict(r) for r in results}
    min_class = int(pd.Series(y_train).value_counts().min()) if len(pd.Series(y_train).value_counts()) else 0
    cv_splits = min(3, min_class) if min_class >= 2 else 2

    for _, row in base_results.iterrows():
        name = row["Model"]
        if name not in param_spaces or name not in trained_models:
            tuning_log.append({"Model": name, "Status": "Skipped", "Reason": "No lightweight tuning space configured"})
            continue
        try:
            base_model = clone(trained_models[name]["model"])
            search = RandomizedSearchCV(
                base_model,
                param_distributions=param_spaces[name],
                n_iter=min(8, sum(len(v) for v in param_spaces[name].values())),
                scoring="accuracy",
                cv=cv_splits,
                random_state=42,
                n_jobs=-1,
            )
            search.fit(X_train, y_train)
            tuned_model = search.best_estimator_
            tuned_preds = tuned_model.predict(X_test)
            tuned_acc = accuracy_score(y_test, tuned_preds)
            tuned_f1 = f1_score(y_test, tuned_preds, average="weighted", zero_division=0)
            old_acc = float(row["Accuracy"])
            if tuned_acc >= old_acc:
                results_by_model[name] = {
                    "Model": name,
                    "Accuracy": tuned_acc,
                    "Weighted F1": tuned_f1,
                    "CV Mean": float(search.best_score_),
                    "CV Std": trained_models[name].get("cv_std", np.nan),
                }
                trained_models[name] = {
                    "model": tuned_model,
                    "predictions": tuned_preds,
                    "accuracy": tuned_acc,
                    "weighted_f1": tuned_f1,
                    "cv_mean": float(search.best_score_),
                    "cv_std": trained_models[name].get("cv_std", np.nan),
                    "best_params": search.best_params_,
                    "tuned": True,
                }
                tuning_log.append({"Model": name, "Status": "Improved/Kept", "Old Accuracy": old_acc, "Tuned Accuracy": tuned_acc, "Best Params": str(search.best_params_)})
            else:
                tuning_log.append({"Model": name, "Status": "Baseline kept", "Old Accuracy": old_acc, "Tuned Accuracy": tuned_acc, "Best Params": str(search.best_params_)})
        except Exception as e:
            tuning_log.append({"Model": name, "Status": "Failed", "Reason": str(e)})

    return list(results_by_model.values()), trained_models, tuning_log


def render_model_comparison_charts(leaderboard):
    """Premium visual model comparison dashboard."""
    if leaderboard is None or len(leaderboard) == 0:
        return
    chart_df = leaderboard.copy().sort_values("Accuracy", ascending=True)

    st.write("### Model Comparison Charts")
    fig_acc, ax_acc = plt.subplots(figsize=(9, 4.5))
    ax_acc.barh(chart_df["Model"], chart_df["Accuracy"] * 100)
    ax_acc.set_xlabel("Accuracy (%)")
    ax_acc.set_title("Accuracy by Model")
    ax_acc.set_xlim(0, 100)
    st.pyplot(fig_acc)

    if "Weighted F1" in chart_df.columns:
        fig_f1, ax_f1 = plt.subplots(figsize=(9, 4.5))
        ax_f1.barh(chart_df["Model"], chart_df["Weighted F1"] * 100)
        ax_f1.set_xlabel("Weighted F1 (%)")
        ax_f1.set_title("Weighted F1 by Model")
        ax_f1.set_xlim(0, 100)
        st.pyplot(fig_f1)

    if "CV Mean" in leaderboard.columns:
        cv_df = leaderboard.dropna(subset=["CV Mean"]).sort_values("CV Mean", ascending=True)
        if len(cv_df) > 0:
            fig_cv, ax_cv = plt.subplots(figsize=(9, 4.5))
            ax_cv.barh(cv_df["Model"], cv_df["CV Mean"] * 100)
            ax_cv.set_xlabel("Cross-validation Mean Accuracy (%)")
            ax_cv.set_title("Cross-validation Stability by Model")
            ax_cv.set_xlim(0, 100)
            st.pyplot(fig_cv)


def render_all_model_roc_curves(trained_models, X_test, y_test):
    """Draw ROC curves for all binary classifiers that support predict_proba."""
    if trained_models is None or len(trained_models) == 0 or len(pd.Series(y_test).unique()) != 2:
        return

    fig, ax = plt.subplots(figsize=(8, 6))
    plotted = 0
    for name, bundle in trained_models.items():
        model = bundle.get("model")
        if not hasattr(model, "predict_proba"):
            continue
        try:
            probs = model.predict_proba(X_test)[:, 1]
            fpr, tpr, _ = roc_curve(y_test, probs)
            roc_auc = auc(fpr, tpr)
            ax.plot(fpr, tpr, label=f"{name} AUC={roc_auc:.3f}")
            plotted += 1
        except Exception:
            continue

    if plotted > 0:
        ax.plot([0, 1], [0, 1], linestyle="--", label="Random baseline")
        ax.set_xlabel("False Positive Rate")
        ax.set_ylabel("True Positive Rate")
        ax.set_title("ROC Curves for All Supported Models")
        ax.legend(loc="lower right", fontsize=8)
        st.pyplot(fig)


def render_model_cards(trained_models, leaderboard):
    """Expandable model cards with strengths, weaknesses, speed notes, and deployment fit."""
    if trained_models is None or leaderboard is None:
        return

    strengths = {
        "Random Forest": "Strong baseline, robust to nonlinear patterns, handles mixed feature effects well.",
        "Extra Trees": "Fast ensemble, often strong on tabular datasets, reduces variance through randomization.",
        "Gradient Boosting": "Powerful boosted trees, good for structured data and feature interactions.",
        "AdaBoost": "Good when weak learners can be combined into a strong classifier; often effective on clean tabular data.",
        "Decision Tree": "Very interpretable and fast, but can overfit.",
        "Logistic Regression": "Fast, stable, and interpretable for linear decision boundaries.",
        "SVM": "Strong margin-based classifier, useful on smaller high-dimensional datasets.",
        "KNN": "Simple non-parametric baseline, useful when similar records have similar labels.",
        "Naive Bayes": "Very fast probabilistic baseline, works well when feature independence is reasonable.",
        "XGBoost": "High-performing gradient boosting model for structured datasets.",
        "LightGBM": "Fast gradient boosting model designed for larger tabular datasets.",
        "CatBoost": "Boosting model known for strong tabular performance and categorical handling.",
    }
    weaknesses = {
        "Decision Tree": "Can overfit and may be unstable across splits.",
        "KNN": "Can slow down with larger datasets and is sensitive to scaling.",
        "SVM": "Can be slower on larger datasets and less directly interpretable.",
        "Naive Bayes": "Feature independence assumption can be unrealistic.",
        "Logistic Regression": "May underfit nonlinear patterns.",
    }

    display = leaderboard.copy().sort_values("Accuracy", ascending=False).reset_index(drop=True)
    for rank, row in display.iterrows():
        name = row["Model"]
        acc = row.get("Accuracy", np.nan)
        f1 = row.get("Weighted F1", np.nan)
        cv_mean = row.get("CV Mean", np.nan)
        cv_std = row.get("CV Std", np.nan)
        bundle = trained_models.get(name, {})
        with st.expander(f"#{rank + 1} {name} — {acc * 100:.2f}% accuracy"):
            c1, c2, c3, c4 = st.columns(4)
            c1.metric("Accuracy", f"{acc * 100:.2f}%" if not pd.isna(acc) else "N/A")
            c2.metric("Weighted F1", f"{f1 * 100:.2f}%" if not pd.isna(f1) else "N/A")
            c3.metric("CV Mean", f"{cv_mean * 100:.2f}%" if not pd.isna(cv_mean) else "N/A")
            c4.metric("CV Std", f"±{cv_std * 100:.2f}%" if not pd.isna(cv_std) else "N/A")
            st.markdown(f"**Strength:** {strengths.get(name, 'Useful baseline for tabular prediction tasks.')}")
            st.markdown(f"**Watch out:** {weaknesses.get(name, 'Validate carefully on external data before deployment.')}")
            st.markdown("**Deployment note:** Use this model only after confirming cross-validation stability and checking feature drivers.")
            if bundle.get("tuned"):
                st.markdown(f"**Tuned parameters:** `{bundle.get('best_params', {})}`")


def render_shap_dashboard(model, X_train, X_test, feature_names):
    """Show SHAP plots when supported, without breaking the app."""
    if shap is None:
        st.info("SHAP is not installed. Run `pip install shap` to unlock SHAP plots.")
        return

    estimator = get_inner_estimator(model)
    if not hasattr(estimator, "feature_importances_"):
        st.info("SHAP TreeExplainer is best supported for tree-based models. This model is using universal explainability instead.")
        return

    try:
        sample = X_test.sample(min(80, len(X_test)), random_state=42) if len(X_test) > 80 else X_test.copy()
        explainer = shap.TreeExplainer(estimator)
        shap_values = explainer.shap_values(sample)
        values_for_plot = shap_values
        if isinstance(shap_values, list) and len(shap_values) > 1:
            values_for_plot = shap_values[1]

        st.write("### SHAP Summary Plot")
        fig_shap = plt.figure(figsize=(10, 5))
        shap.summary_plot(values_for_plot, sample, feature_names=feature_names, show=False, plot_type="bar")
        st.pyplot(fig_shap, clear_figure=True)

        st.write("### SHAP Beeswarm Plot")
        fig_bee = plt.figure(figsize=(10, 5))
        shap.summary_plot(values_for_plot, sample, feature_names=feature_names, show=False)
        st.pyplot(fig_bee, clear_figure=True)
    except Exception as e:
        st.info(f"SHAP plot could not be rendered for this model, so LabMind is using universal explainability. Details: {e}")

def render_labmind_logo(size="large"):
    """Render the LabMind logo and brand text."""
    import base64
    from pathlib import Path

    size_class = "brand-large" if size == "large" else "brand-small"

    logo_path = Path(__file__).parent / "labmind_logo.png"

    if not logo_path.exists():
        st.error(f"Logo not found at: {logo_path}")
        return

    logo_base64 = base64.b64encode(
        logo_path.read_bytes()
    ).decode("utf-8")

    st.markdown(
        f"""
        <div class="brand-lockup {size_class}">
            <img
                src="data:image/png;base64,{logo_base64}"
                class="brand-logo"
                alt="LabMind.ai logo"
            >
            <div>
                <div class="brand-name">LabMind.ai</div>
                <div class="brand-tagline">AI Data Analyst</div>
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )

def render_pricing_page():

    selected_plan = st.session_state.get("selected_pricing_plan")

    if selected_plan == "professional":
        st.markdown("## 🚀 Upgrade to Professional")
        st.markdown(
            """
            ### Professional Plan — $29/month

            Unlock the full LabMind experience:

            - ✅ Expanded dataset usage
            - ✅ Full AutoML model suite
            - ✅ AI Insights
            - ✅ Explainable AI
            - ✅ Prediction Playground
            - ✅ PDF & PowerPoint reports
            - ✅ 500 MB uploads
            - ✅ Email support
            """
        )

        if st.button("← Back to Pricing", key="back_from_pro"):
            st.session_state["selected_pricing_plan"] = None
            st.rerun()

        st.button(
            "Continue to Payment — $29/month",
            type="primary",
            use_container_width=True,
            key="continue_pro_payment",
        )

        return

    if selected_plan == "business":
        st.markdown("## 🏢 Upgrade to Business")
        st.markdown(
            """
            ### Business Plan — $99/month

            Everything in Professional, plus:

            - ✅ Higher usage limits
            - ✅ Priority processing
            - ✅ Larger dataset uploads
            - ✅ Team workspaces — coming soon
            - ✅ Project sharing — coming soon
            - ✅ API access — planned
            - ✅ Priority support
            - ✅ Early access to new features
            """
        )

        if st.button("← Back to Pricing", key="back_from_business"):
            st.session_state["selected_pricing_plan"] = None
            st.rerun()

        st.button(
            "Continue to Payment — $99/month",
            type="primary",
            use_container_width=True,
            key="continue_business_payment",
        )

        return

    if selected_plan == "free":
        st.markdown("## 🎉 Free Plan")
        st.success("You're ready to use LabMind on the Free plan.")

        st.markdown(
            """
            Your Free plan includes:

            - ✅ 3 datasets per month
            - ✅ Basic dataset analysis
            - ✅ Basic cleaning
            - ✅ Limited AutoML
            - ✅ Basic visualizations
            - ✅ 25 MB uploads
            """
        )

        if st.button("← Back to Pricing", key="back_from_free"):
            st.session_state["selected_pricing_plan"] = None
            st.rerun()

        return
    st.markdown(
        """
        <style>
        .pricing-wrap {
            max-width: 1200px;
            margin: 0 auto;
            padding: 24px 0 50px 0;
        }

        .pricing-title {
            text-align: center;
            color: #FFFFFF;
            font-size: 46px;
            font-weight: 900;
            margin-bottom: 8px;
        }

        .pricing-subtitle {
            text-align: center;
            color: #94A3B8;
            font-size: 18px;
            margin-bottom: 38px;
        }

        .pricing-card {
            min-height: 540px;
            padding: 28px;
            border-radius: 26px;
            background: linear-gradient(
                160deg,
                rgba(30, 41, 59, 0.98),
                rgba(15, 23, 42, 0.98)
            );
            border: 1px solid rgba(148, 163, 184, 0.18);
            box-shadow: 0 24px 60px rgba(2, 6, 23, 0.28);
            position: relative;
        }

        .pricing-card-popular {
            border: 1px solid rgba(99, 102, 241, 0.65);
            box-shadow: 0 26px 70px rgba(79, 70, 229, 0.24);
            transform: translateY(-8px);
        }

        .pricing-badge {
            display: inline-block;
            padding: 7px 12px;
            border-radius: 999px;
            color: #E0E7FF;
            background: rgba(99, 102, 241, 0.18);
            border: 1px solid rgba(99, 102, 241, 0.35);
            font-size: 12px;
            font-weight: 800;
            margin-bottom: 16px;
        }

        .pricing-plan {
            color: #FFFFFF;
            font-size: 24px;
            font-weight: 900;
            margin-bottom: 4px;
        }

        .pricing-price {
            color: #FFFFFF;
            font-size: 50px;
            font-weight: 1000;
            line-height: 1;
            margin-top: 18px;
        }

        .pricing-price span {
            color: #94A3B8;
            font-size: 16px;
            font-weight: 700;
        }

        .pricing-desc {
            color: #CBD5E1;
            font-size: 14px;
            line-height: 1.7;
            margin: 18px 0 24px 0;
        }

        .pricing-feature {
            color: #E2E8F0;
            font-size: 14px;
            margin-bottom: 12px;
        }

        .pricing-feature strong {
            color: #FFFFFF;
        }

        .pricing-note {
            text-align: center;
            color: #64748B;
            font-size: 13px;
            margin-top: 28px;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )

    st.markdown(
        """
        <div class="pricing-wrap">
            <div class="pricing-title">Choose the plan that fits your workflow</div>
            <div class="pricing-subtitle">
                Start free. Upgrade when you need more power, automation, and scale.
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )

    free_col, pro_col, business_col = st.columns(3, gap="large")

    with free_col:
        free_html = """
        <div class="pricing-card">
            <div class="pricing-plan">Free</div>
            <div class="pricing-price">$0<span>/month</span></div>
            <div class="pricing-desc">
                Perfect for exploring LabMind and trying core dataset workflows.
            </div>
            <div class="pricing-feature">✓ 3 datasets per month</div>
            <div class="pricing-feature">✓ Basic dataset analysis</div>
            <div class="pricing-feature">✓ Basic cleaning tools</div>
            <div class="pricing-feature">✓ Limited AutoML</div>
            <div class="pricing-feature">✓ Basic visualizations</div>
            <div class="pricing-feature">✓ 25 MB upload limit</div>
        </div>
        """

        free_html = "".join(
            line.strip() for line in free_html.splitlines()
        )

        st.markdown(
            free_html,
            unsafe_allow_html=True,
        )

        if st.button(
            "Get Started",
            use_container_width=True,
            key="pricing_free",
        ):
            st.session_state["selected_pricing_plan"] = "free"
            st.rerun()


    with pro_col:
        pro_html = """
        <div class="pricing-card pricing-card-popular">
            <div class="pricing-badge">MOST POPULAR</div>
            <div class="pricing-plan">Professional</div>
            <div class="pricing-price">$29<span>/month</span></div>
            <div class="pricing-desc">
                For researchers, analysts, freelancers, and power users.
            </div>
            <div class="pricing-feature">✓ Everything in Free</div>
            <div class="pricing-feature">✓ Expanded dataset usage</div>
            <div class="pricing-feature">✓ Full AutoML model suite</div>
            <div class="pricing-feature">✓ AI Insights</div>
            <div class="pricing-feature">✓ Explainable AI</div>
            <div class="pricing-feature">✓ Prediction Playground</div>
            <div class="pricing-feature">✓ PDF & PPT reports</div>
            <div class="pricing-feature">✓ 500 MB upload limit</div>
            <div class="pricing-feature">✓ Email support</div>
        </div>
        """

        pro_html = "".join(
            line.strip() for line in pro_html.splitlines()
        )

        st.markdown(
            pro_html,
            unsafe_allow_html=True,
        )

        if st.button(
            "Upgrade to Professional",
            use_container_width=True,
            type="primary",
            key="pricing_pro",
        ):
            st.session_state["selected_pricing_plan"] = "professional"
            st.rerun()


    with business_col:
        business_html = """
        <div class="pricing-card">
            <div class="pricing-plan">Business</div>
            <div class="pricing-price">$99<span>/month</span></div>
            <div class="pricing-desc">
                Built for teams and organizations running serious analytics workflows.
            </div>
            <div class="pricing-feature">✓ Everything in Professional</div>
            <div class="pricing-feature">✓ Higher usage limits</div>
            <div class="pricing-feature">✓ Priority processing</div>
            <div class="pricing-feature">✓ Larger dataset uploads</div>
            <div class="pricing-feature">✓ Team workspaces <strong>(coming soon)</strong></div>
            <div class="pricing-feature">✓ Project sharing <strong>(coming soon)</strong></div>
            <div class="pricing-feature">✓ API access <strong>(planned)</strong></div>
            <div class="pricing-feature">✓ Priority support</div>
            <div class="pricing-feature">✓ Early access to new features</div>
        </div>
        """

        business_html = "".join(
            line.strip() for line in business_html.splitlines()
        )

        st.markdown(
            business_html,
            unsafe_allow_html=True,
        )

        if st.button(
            "Choose Business",
            use_container_width=True,
            key="pricing_business",
        ):
            st.session_state["selected_pricing_plan"] = "business"
            st.rerun()

    st.markdown(
        """
        <div class="pricing-note">
            Annual billing and enterprise custom plans will be available soon.
        </div>
        """,
        unsafe_allow_html=True,
    )

def infer_dataset_difficulty(df, numeric_cols, text_cols, missing_percentage):
    """A simple product-style difficulty estimate for users."""
    score = 0
    if df.shape[0] > 5000:
        score += 2
    elif df.shape[0] > 1000:
        score += 1
    if df.shape[1] > 50:
        score += 2
    elif df.shape[1] > 20:
        score += 1
    if missing_percentage > 20:
        score += 2
    elif missing_percentage > 5:
        score += 1
    if len(text_cols) > len(numeric_cols):
        score += 1
    if score >= 4:
        return "Hard"
    if score >= 2:
        return "Medium"
    return "Easy"


def render_ai_dataset_summary(df, numeric_cols, text_cols, health_score, missing_percentage, duplicate_rows, default_target_name):
    """ChatGPT-style instant dataset summary."""
    difficulty = infer_dataset_difficulty(df, numeric_cols, text_cols, missing_percentage)
    target_text = default_target_name if default_target_name else "not selected yet"
    algo = "Random Forest / Extra Trees" if len(numeric_cols) >= 3 else "Logistic Regression baseline"
    train_time = "< 5 seconds" if df.shape[0] < 5000 else "under a minute for baseline models"

    st.markdown(f"""
    <div class="ai-summary-card">
        <div class="ai-summary-title">✨ AI Dataset Summary</div>
        <div class="summary-grid">
            <div><b>{df.shape[0]:,}</b><span>Rows detected</span></div>
            <div><b>{df.shape[1]:,}</b><span>Columns detected</span></div>
            <div><b>{len(numeric_cols)}</b><span>Numeric features</span></div>
            <div><b>{len(text_cols)}</b><span>Text/category features</span></div>
        </div>
        <div class="summary-text">
            LabMind thinks the likely target is <b>{target_text}</b>. Dataset health is <b>{health_score}/100</b>,
            missing data is <b>{missing_percentage:.2f}%</b>, duplicates detected: <b>{duplicate_rows}</b>.
            Recommended baseline: <b>{algo}</b>. Difficulty: <b>{difficulty}</b>. Estimated baseline training time: <b>{train_time}</b>.
        </div>
    </div>
    """, unsafe_allow_html=True)


def render_premium_recommendations(df, numeric_cols, text_cols, recommendations):
    """Higher-level product recommendations."""
    likely_id_cols = [c for c in df.columns if str(c).lower() in ["id", "passengerid", "patient_id"] or str(c).lower().endswith("id")]
    items = []
    if likely_id_cols:
        items.append(f"Avoid using ID-like columns as predictors: {', '.join(map(str, likely_id_cols[:5]))}.")
    if len(numeric_cols) > 0:
        items.append("Scale numerical features for distance/linear models such as SVM, KNN, and Logistic Regression.")
    if len(text_cols) > 0:
        items.append("Encode categorical/text columns safely before model training.")
    items.append("Use cross-validation and model leaderboard results before choosing a deployment model.")
    items.extend(recommendations[:3])

    html = "".join([f"<li>✅ {x}</li>" for x in items])
    st.markdown(f"""
    <div class="ai-summary-card compact-card">
        <div class="ai-summary-title">🧠 LabMind Recommendations</div>
        <ul class="premium-list">{html}</ul>
    </div>
    """, unsafe_allow_html=True)


def render_model_comparison_dashboard(leaderboard):
    """Premium non-table model comparison blocks."""
    if leaderboard is None or len(leaderboard) == 0:
        return

    top = (
        leaderboard
        .sort_values(
            ["Accuracy", "Weighted F1", "CV Mean"],
            ascending=[False, False, False],
            na_position="last",
        )
        .head(6)
    )

    rows = []

    for _, r in top.iterrows():
        acc = float(r.get("Accuracy", 0)) * 100

        weighted_f1 = r.get("Weighted F1", np.nan)
        f1 = (
            float(weighted_f1) * 100
            if not pd.isna(weighted_f1)
            else acc
        )

        if r["Model"] in ["Decision Tree", "Logistic Regression"]:
            interpretability = 95
        elif r["Model"] in [
            "Random Forest",
            "Extra Trees",
            "AdaBoost",
            "Gradient Boosting",
            "XGBoost",
            "LightGBM",
            "CatBoost",
        ]:
            interpretability = 75
        else:
            interpretability = 60

        deploy = min(98, max(40, (acc + f1) / 2))

        rows.append(
            f'<div class="model-score-card">'
            f'<div class="model-score-title">{r["Model"]}</div>'
            f'<div class="score-line"><span>Accuracy</span>'
            f'<b>{acc:.2f}%</b></div>'
            f'<div class="bar"><i style="width:{acc:.1f}%"></i></div>'
            f'<div class="score-line"><span>Weighted F1</span>'
            f'<b>{f1:.2f}%</b></div>'
            f'<div class="bar"><i style="width:{f1:.1f}%"></i></div>'
            f'<div class="score-line"><span>Interpretability</span>'
            f'<b>{interpretability}%</b></div>'
            f'<div class="bar"><i style="width:{interpretability}%"></i></div>'
            f'<div class="score-line"><span>Deployment Fit</span>'
            f'<b>{deploy:.0f}%</b></div>'
            f'<div class="bar"><i style="width:{deploy:.1f}%"></i></div>'
            f'</div>'
        )

    comparison_html = (
        '<h3 class="premium-heading">📊 Premium Model Comparison</h3>'
        '<div class="model-score-grid">'
        + "".join(rows)
        + "</div>"
    )

    st.markdown(comparison_html, unsafe_allow_html=True)

def render_ai_model_card(best_model_name, best_accuracy, leaderboard=None):
    """Executive model card for the winning model."""
    row = None
    if leaderboard is not None and len(leaderboard) > 0:
        row = leaderboard[leaderboard["Model"] == best_model_name]
        row = row.iloc[0] if len(row) else None
    cv = row.get("CV Mean", np.nan) if row is not None else np.nan
    cv_text = "N/A" if pd.isna(cv) else f"{cv * 100:.2f}%"
    risk = "Low" if best_accuracy >= 0.9 else "Medium" if best_accuracy >= 0.75 else "High"
    stars = "⭐⭐⭐⭐⭐" if best_accuracy >= 0.9 else "⭐⭐⭐⭐" if best_accuracy >= 0.75 else "⭐⭐⭐"
    st.markdown(f"""
    <div class="winner-card">
        <div class="winner-title">🏆 {best_model_name}</div>
        <div class="summary-grid">
            <div><b>{best_accuracy * 100:.2f}%</b><span>Validation Accuracy</span></div>
            <div><b>{cv_text}</b><span>Cross-validation</span></div>
            <div><b>{risk}</b><span>Overfitting Risk</span></div>
            <div><b>{stars}</b><span>Recommendation</span></div>
        </div>
        <p>Use this as the current baseline model. Before real deployment, validate it on fresh data and review explainability results.</p>
    </div>
    """, unsafe_allow_html=True)


def render_feature_impact_blocks(explain_df, top_n=8):
    """Render horizontal feature-impact bars without Markdown escaping."""
    if explain_df is None or len(explain_df) == 0:
        return

    clean = (
        explain_df.copy()
        .sort_values("Importance", ascending=False)
        .head(top_n)
    )

    clean["Importance"] = pd.to_numeric(
        clean["Importance"],
        errors="coerce",
    ).fillna(0)

    total = float(clean["Importance"].sum())

    if total <= 0:
        st.info("No positive feature importance values available.")
        return

    rows = []

    for _, row in clean.iterrows():
        pct = float(row["Importance"]) / total * 100
        feature_name = str(row["Feature"])

        row_html = (
            '<div class="feature-impact-row">'
            f'<div class="feature-impact-name">{feature_name}</div>'
            '<div class="feature-impact-bar">'
            f'<i style="width:{min(100, pct):.1f}%"></i>'
            '</div>'
            f'<div class="feature-impact-pct">{pct:.2f}%</div>'
            '</div>'
        )

        rows.append(row_html)

    html = (
        '<h3 class="premium-heading">🔥 Feature Impact Breakdown</h3>'
        '<div class="feature-impact-card">'
        + "".join(rows)
        + "</div>"
    )

    st.markdown(html, unsafe_allow_html=True)

def render_prediction_reason_and_similar(input_df, decoded_prediction, confidence_text):
    """Explain an individual prediction and show closest similar samples."""
    explain_df = st.session_state.get("explain_df")
    top_features = []
    if explain_df is not None and len(explain_df) > 0:
        top_features = explain_df.head(5)["Feature"].astype(str).tolist()
    reason_lines = "".join([f"<li>{f} strongly influenced this model.</li>" for f in top_features]) or "<li>Train explainability to see the strongest drivers.</li>"
    st.markdown(f"""
    <div class="winner-card">
        <div class="winner-title">🎯 Prediction Explanation</div>
        <p>LabMind predicts <b>{decoded_prediction}</b> with confidence <b>{confidence_text}</b>.</p>
        <p><b>Main reasons LabMind checked:</b></p>
        <ul class="premium-list">{reason_lines}</ul>
    </div>
    """, unsafe_allow_html=True)

    try:
        X_train = st.session_state.get("X_train")
        y_train = st.session_state.get("y_train")
        target_label_map = st.session_state.get("target_label_map")
        if X_train is not None and y_train is not None and len(X_train) > 0:
            aligned = input_df[X_train.columns]
            scale = X_train.std(numeric_only=True).replace(0, 1).fillna(1)
            diffs = ((X_train - aligned.iloc[0]) / scale).abs().sum(axis=1).sort_values().head(5)
            similar = X_train.loc[diffs.index].copy()
            similar["Actual Target"] = [decode_target_value(v, target_label_map) for v in y_train.loc[diffs.index]]
            similar["Similarity Rank"] = range(1, len(similar) + 1)
            st.write("### Closest Similar Samples")
            st.dataframe(similar[["Similarity Rank", "Actual Target"] + list(X_train.columns[:8])], use_container_width=True)
    except Exception:
        pass



st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800;900&display=swap');

html, body, [class*="css"] {
    font-family: 'Inter', sans-serif;
}

button,
input,
select,
textarea {
    font-family: 'Inter', sans-serif !important;
}
.stApp {
    background: linear-gradient(135deg, #020617 0%, #0F172A 45%, #111827 100%);
    color: #F8FAFC;
}

.block-container {
    padding-top: 3rem;
    padding-bottom: 3rem;
    max-width: 1250px;
}

section[data-testid="stSidebar"] {
    background: linear-gradient(180deg, #F8FAFC 0%, #E2E8F0 100%);
    border-right: 1px solid #CBD5E1;
}

section[data-testid="stSidebar"] * {
    color: #0F172A !important;
}

.sidebar-title {
    font-size: 28px;
    font-weight: 900;
    margin-bottom: 8px;
}

.sidebar-subtitle {
    font-size: 15px;
    color: #475569 !important;
    margin-bottom: 28px;
}

.sidebar-item {
    padding: 12px 14px;
    border-radius: 12px;
    margin-bottom: 8px;
    font-weight: 800;
}

.hero-box {
    background: radial-gradient(circle at top left, #60A5FA 0%, #2563EB 35%, #7C3AED 100%);
    padding: 44px;
    border-radius: 30px;
    margin-bottom: 30px;
    box-shadow: 0 24px 80px rgba(37, 99, 235, 0.35);
    border: 1px solid rgba(255, 255, 255, 0.22);
}

.main-title {
    font-size: 58px;
    font-weight: 900;
    color: #FFFFFF;
    letter-spacing: -1px;
    margin-bottom: 8px;
}

.subtitle {
    font-size: 24px;
    color: #DBEAFE;
    margin-bottom: 18px;
    font-weight: 700;
}

.hero-text {
    font-size: 18px;
    line-height: 1.7;
    color: #EEF2FF;
}

.upload-wrapper {
    background: linear-gradient(180deg, #1E293B, #111827);
    border-radius: 24px;
    border: 1px solid rgba(147, 197, 253, 0.35);
    padding: 28px;
    margin-bottom: 20px;
}

.upload-title {
    color: white;
    font-size: 26px;
    font-weight: 900;
    margin-bottom: 8px;
}

.upload-subtitle {
    color: #E2E8F0;
    font-size: 16px;
}

[data-testid="stFileUploader"] {
    background: linear-gradient(135deg, #1E293B, #0F172A) !important;
    border: 3px dashed #60A5FA !important;
    border-radius: 24px !important;
    padding: 28px !important;
}

[data-testid="stFileUploader"] label {
    color: white !important;
    font-size: 18px !important;
    font-weight: 900 !important;
}

[data-testid="stFileUploader"] button {
    background: linear-gradient(135deg, #2563EB, #7C3AED) !important;
    color: white !important;
    border-radius: 16px !important;
    font-weight: 900 !important;
    border: none !important;
}

.card {
    background: linear-gradient(180deg, #1E293B, #111827);
    padding: 24px;
    border-radius: 22px;
    border: 1px solid rgba(147, 197, 253, 0.35);
    min-height: 120px;
    box-shadow: 0 16px 40px rgba(0,0,0,0.32);
}

.card-title {
    font-size: 17px;
    color: #CBD5E1;
    font-weight: 800;
    margin-bottom: 10px;
}

.card-value {
    font-size: 34px;
    font-weight: 900;
    color: white;
    word-break: break-word;
}

.section-title {
    font-size: 34px;
    font-weight: 900;
    color: white;
    margin-top: 28px;
    margin-bottom: 16px;
}

.recommendation-box {
    background: linear-gradient(180deg, #1E293B, #111827);
    border-left: 5px solid #3B82F6;
    padding: 18px;
    border-radius: 16px;
    margin-bottom: 14px;
    border-top: 1px solid rgba(147, 197, 253, 0.22);
    border-right: 1px solid rgba(147, 197, 253, 0.22);
    border-bottom: 1px solid rgba(147, 197, 253, 0.22);
}

.recommendation-title {
    font-weight: 900;
    color: white;
    font-size: 18px;
    margin-bottom: 6px;
}

.recommendation-text {
    color: #E2E8F0;
    font-size: 15px;
    line-height: 1.6;
}

.faq-grid {
    display: grid;
    grid-template-columns: repeat(2, minmax(0, 1fr));
    gap: 14px;
    margin-top: 14px;
    margin-bottom: 22px;
}

.faq-card {
    background: linear-gradient(180deg, #1E293B, #111827);
    border: 1px solid rgba(147, 197, 253, 0.30);
    border-radius: 16px;
    padding: 16px;
    color: #E2E8F0;
    font-weight: 700;
}

.chat-box {
    background: linear-gradient(180deg, #1E293B, #111827);
    border-radius: 20px;
    padding: 22px;
    border: 1px solid rgba(147, 197, 253, 0.35);
    margin-top: 16px;
}

.chat-answer {
    color: #E2E8F0;
    font-size: 16px;
    line-height: 1.7;
    white-space: pre-wrap;
}

.stTabs [data-baseweb="tab-list"] {
    gap: 10px;
    border-bottom: 1px solid rgba(148, 163, 184, 0.18);
}

.stTabs [data-baseweb="tab"] {
    background-color: rgba(30, 41, 59, 0.95);
    border-radius: 14px 14px 0 0;
    padding: 14px 20px;
    color: #CBD5E1;
    font-weight: 800;
}

.stTabs [aria-selected="true"] {
    background: linear-gradient(135deg, #2563EB, #7C3AED);
    color: white !important;
}

.stButton > button {
    background: linear-gradient(135deg, #2563EB, #7C3AED);
    color: white;
    border: none;
    border-radius: 16px;
    padding: 14px 30px;
    font-weight: 900;
    font-size: 16px;
}

.stDownloadButton > button {
    background: linear-gradient(135deg, #16A34A, #22C55E);
    color: white;
    border: none;
    border-radius: 16px;
    padding: 14px 30px;
    font-weight: 900;
    font-size: 16px;
}

input, textarea {
    color: #0F172A !important;
    background: #F8FAFC !important;
    border-radius: 14px !important;
}
            
/* Fix Streamlit metric text on dark background */
[data-testid="stMetricLabel"],
[data-testid="stMetricValue"],
[data-testid="stMetricDelta"] {
    color: #FFFFFF !important;
}

[data-testid="stMetricValue"] {
    font-size: 34px !important;
    font-weight: 900 !important;
    text-shadow: 0 0 18px rgba(255,255,255,0.25);
}

[data-testid="stMetricLabel"] {
    color: #CBD5E1 !important;
    font-weight: 800 !important;
}
                        
.insight-card {
    background: linear-gradient(180deg, #10233F, #0F172A);
    border: 1px solid rgba(96, 165, 250, 0.40);
    border-left: 6px solid #60A5FA;
    border-radius: 22px;
    padding: 26px;
    margin: 18px 0 24px 0;
    box-shadow: 0 18px 45px rgba(0,0,0,0.28);
}

.insight-card h3 {
    color: #FFFFFF !important;
    margin-top: 0;
    font-size: 26px;
    font-weight: 900;
}

.insight-card p {
    color: #F8FAFC !important;
    font-size: 17px;
    line-height: 1.75;
}

.insight-card li {
    color: #F8FAFC !important;
    font-size: 17px;
    line-height: 1.65;
}

.model-highlight {
    background: linear-gradient(135deg, rgba(37, 99, 235, 0.28), rgba(124, 58, 237, 0.28));
    border: 1px solid rgba(147, 197, 253, 0.35);
    border-radius: 18px;
    padding: 18px;
    margin: 16px 0;
    color: #FFFFFF !important;
    font-size: 17px;
    line-height: 1.7;
}



/* LabMind V3 brand polish */
.brand-lockup {
    display: flex;
    align-items: center;
    gap: 14px;
    margin-bottom: 18px;
}

.brand-logo {
    width: 58px;
    height: 58px;
    object-fit: contain;
    flex-shrink: 0;
}

.brand-name {
    font-size: 42px;
    font-weight: 1000;
    color: #FFFFFF;
    letter-spacing: -1px;
}

.brand-tagline {
    font-size: 16px;
    color: #BDEAFE;
    font-weight: 800;
}
.brand-small .brand-logo {
    width: 42px;
    height: 42px;
}.brand-small .brand-name {font-size:27px; color:#0F172A;}
.brand-small .brand-tagline {color:#475569;}
.labmind-hero {
    background:
        radial-gradient(
            circle at top left,
            rgba(56, 189, 248, 0.28),
            transparent 34%
        ),
        linear-gradient(
            135deg,
            #2563EB 0%,
            #4F46E5 55%,
            #7C3AED 100%
        );
    border: 1px solid rgba(255, 255, 255, 0.20);
    border-radius: 30px;
    padding: 42px 48px;
    margin: 12px 0 32px 0;
    box-shadow: 0 24px 70px rgba(37, 99, 235, 0.25);
    overflow: hidden;
}

.hero-brand-row {
    display: flex;
    align-items: center;
    gap: 18px;
    margin-bottom: 32px;
}

.hero-brand-logo {
    width: 72px;
    height: 72px;
    object-fit: contain;
    flex-shrink: 0;
    background: rgba(255, 255, 255, 0.96);
    border-radius: 18px;
    padding: 8px;
    box-shadow: 0 14px 32px rgba(15, 23, 42, 0.22);
}

.hero-logo-fallback {
    width: 72px;
    height: 72px;
    display: flex;
    align-items: center;
    justify-content: center;
    background: rgba(255, 255, 255, 0.96);
    border-radius: 18px;
    font-size: 38px;
}

.hero-brand-name {
    color: #FFFFFF;
    font-size: 44px;
    font-weight: 1000;
    letter-spacing: -1.5px;
    line-height: 1;
}

.hero-brand-tagline {
    color: #DBEAFE;
    font-size: 16px;
    font-weight: 800;
    margin-top: 10px;
}

.hero-brand-text-row {
    margin-bottom: 28px;
}

.hero-brand-name-new {
    color: #FFFFFF;
    font-size: 46px;
    font-weight: 1000;
    letter-spacing: -1.5px;
    line-height: 1;
}

.hero-brand-tagline-new {
    color: #DBEAFE;
    font-size: 16px;
    font-weight: 800;
    margin-top: 10px;
}

.hero-main-title {
    color: #FFFFFF;
    font-size: 40px;
    font-weight: 900;
    line-height: 1.15;
    letter-spacing: -1.2px;
    margin-bottom: 18px;
}

.hero-title-accent {
    display: inline-block;
    color: #DBEAFE;
    margin-top: 6px;
}

.hero-description {
    color: #EFF6FF;
    font-size: 18px;
    line-height: 1.75;
    max-width: 980px;
}

.hero-feature-row {
    display: flex;
    flex-wrap: wrap;
    gap: 12px;
    margin-top: 28px;
}

.hero-feature-row span {
    color: #FFFFFF;
    font-size: 14px;
    font-weight: 800;
    background: rgba(15, 23, 42, 0.20);
    border: 1px solid rgba(255, 255, 255, 0.20);
    border-radius: 999px;
    padding: 9px 14px;
    backdrop-filter: blur(8px);
}

.hero-feature-row span {
    transition: transform 0.2s ease, border-color 0.2s ease,
                background 0.2s ease, box-shadow 0.2s ease;
}

.hero-feature-row span:hover {
    transform: translateY(-3px);
    border-color: rgba(255, 255, 255, 0.55);
    background: rgba(15, 23, 42, 0.32);
    box-shadow: 0 10px 24px rgba(15, 23, 42, 0.24);
}

.hero-brand-inline {
    display: flex;
    align-items: center;
    gap: 16px;
    margin-bottom: 28px;
}

.hero-brand-logo-box {
    width: 68px;
    height: 68px;
    border-radius: 18px;
    background: rgba(255, 255, 255, 0.96);
    display: flex;
    align-items: center;
    justify-content: center;
    font-size: 34px;
    box-shadow: 0 14px 32px rgba(15, 23, 42, 0.22);
    flex-shrink: 0;
}

.hero-real-logo {
    width: 46px;
    height: 46px;
    object-fit: contain;
    display: block;
}

@media (max-width: 800px) {
    .labmind-hero {
        padding: 30px 24px;
        border-radius: 24px;
    }

    .hero-brand-logo,
    .hero-logo-fallback {
        width: 58px;
        height: 58px;
    }

    .hero-brand-name {
        font-size: 34px;
    }

    .hero-main-title {
        font-size: 26px;
    }

    .hero-description {
        font-size: 16px;
    }
 
    }

    /* ========================
Premium Dataset Health Card
======================== */

.dataset-health-card {
    position: relative;
    overflow: hidden;
    margin: 22px 0 30px 0;
    padding: 28px;
    border-radius: 24px;
    background:
        linear-gradient(
            145deg,
            rgba(30, 41, 59, 0.98),
            rgba(15, 23, 42, 0.98)
        );
    border: 1px solid rgba(148, 163, 184, 0.22);
    box-shadow: 0 20px 50px rgba(2, 6, 23, 0.28);
}

.dataset-health-card::before {
    content: "";
    position: absolute;
    width: 190px;
    height: 190px;
    top: -85px;
    right: -65px;
    border-radius: 50%;
    background: rgba(99, 102, 241, 0.18);
    filter: blur(2px);
    pointer-events: none;
}

.dataset-health-header {
    position: relative;
    z-index: 1;
    display: flex;
    align-items: center;
    justify-content: space-between;
    gap: 24px;
    margin-bottom: 18px;
}

.dataset-health-eyebrow {
    margin-bottom: 7px;
    color: #93C5FD;
    font-size: 13px;
    font-weight: 900;
    letter-spacing: 1.4px;
    text-transform: uppercase;
}

.dataset-health-title {
    color: #FFFFFF;
    font-size: 25px;
    font-weight: 900;
    line-height: 1.25;
}

.dataset-health-score-wrap {
    display: flex;
    align-items: baseline;
    justify-content: center;
    min-width: 132px;
    padding: 14px 18px;
    border-radius: 18px;
    background: linear-gradient(
        135deg,
        rgba(59, 130, 246, 0.28),
        rgba(124, 58, 237, 0.32)
    );
    border: 1px solid rgba(147, 197, 253, 0.32);
}

.dataset-health-score {
    color: #FFFFFF;
    font-size: 42px;
    font-weight: 1000;
    line-height: 1;
}

.dataset-health-total {
    margin-left: 4px;
    color: #BFDBFE;
    font-size: 16px;
    font-weight: 800;
}

.dataset-health-status {
    display: inline-flex;
    align-items: center;
    margin-bottom: 22px;
    padding: 8px 14px;
    border-radius: 999px;
    font-size: 14px;
    font-weight: 900;
}

.health-excellent {
    color: #BBF7D0;
    background: rgba(34, 197, 94, 0.17);
    border: 1px solid rgba(34, 197, 94, 0.35);
}

.health-good {
    color: #BFDBFE;
    background: rgba(59, 130, 246, 0.17);
    border: 1px solid rgba(59, 130, 246, 0.35);
}

.health-warning {
    color: #FDE68A;
    background: rgba(245, 158, 11, 0.17);
    border: 1px solid rgba(245, 158, 11, 0.35);
}

.health-metrics-grid {
    position: relative;
    z-index: 1;
    display: grid;
    grid-template-columns: repeat(2, minmax(0, 1fr));
    gap: 16px;
}

.health-metric {
    padding: 17px;
    border-radius: 17px;
    background: rgba(15, 23, 42, 0.52);
    border: 1px solid rgba(148, 163, 184, 0.18);
}

.health-metric-row {
    display: flex;
    align-items: center;
    justify-content: space-between;
    gap: 12px;
    margin-bottom: 11px;
}

.health-metric-row span {
    color: #CBD5E1;
    font-size: 14px;
    font-weight: 800;
}

.health-metric-row strong {
    color: #FFFFFF;
    font-size: 17px;
    font-weight: 1000;
}

.health-progress {
    width: 100%;
    height: 9px;
    overflow: hidden;
    border-radius: 999px;
    background: rgba(148, 163, 184, 0.20);
}

.health-progress > div {
    height: 100%;
    border-radius: 999px;
    background: linear-gradient(
        90deg,
        #38BDF8,
        #6366F1,
        #A855F7
    );
    box-shadow: 0 0 18px rgba(99, 102, 241, 0.50);
}

@media (max-width: 800px) {
    .dataset-health-header {
        align-items: flex-start;
        flex-direction: column;
    }

    .dataset-health-score-wrap {
        min-width: 112px;
    }

    .health-metrics-grid {
        grid-template-columns: 1fr;
    }
}

    /* =========================
   AI Workspace Stats
   ========================= */

.workspace-stats-grid {
    display: grid;
    grid-template-columns: repeat(4, minmax(0, 1fr));
    gap: 18px;
    margin: 0 0 28px 0;
}

.workspace-stat-card{
    position:relative;
    overflow:hidden;

    background:linear-gradient(
        145deg,
        rgba(34,52,90,.98),
        rgba(20,28,45,.98)
    );

    border:1px solid rgba(255,255,255,.08);

    border-radius:24px;

    padding:26px;

    min-height:185px;

    transition:.3s ease;

    box-shadow:
        0 10px 30px rgba(0,0,0,.25);
}

.workspace-stat-card:hover{

    transform:translateY(-8px);

    border-color:#6C7DFF;

    box-shadow:
        0 25px 60px rgba(74,98,255,.25);
}

.workspace-stat-card::before {
    content: "";
    position: absolute;
    width: 110px;
    height: 110px;
    right: -35px;
    top: -35px;
    border-radius: 50%;
    background: rgba(99, 102, 241, 0.18);
    filter: blur(2px);
}

.workspace-stat-card:hover {
    transform: translateY(-5px);
    border-color: rgba(96, 165, 250, 0.55);
    box-shadow: 0 24px 58px rgba(37, 99, 235, 0.22);
}

.workspace-stat-top{
    display:flex;
    justify-content:space-between;
    align-items:center;
    margin-bottom:22px;
}

.workspace-stat-icon{
    width:60px;
    height:60px;

    display:flex;
    align-items:center;
    justify-content:center;

    border-radius:18px;

    font-size:30px;

    background:linear-gradient(
        135deg,
        rgba(91,125,255,.22),
        rgba(143,92,246,.22)
    );

    color:white;
}

.workspace-stat-value{
    font-size:54px;
    font-weight:800;
    color:white;
    margin-top:18px;
    line-height:1;
}

.workspace-stat-value {
    color: #FFFFFF;
    font-size: 32px;
    font-weight: 900;
    letter-spacing: -1px;
    line-height: 1;
}

.workspace-stat-caption{
    margin-top:14px;
    color:#C9D3E8;
    font-size:15px;
}

@media (max-width: 1050px) {
    .workspace-stats-grid {
        grid-template-columns: repeat(2, minmax(0, 1fr));
    }
}

@media (max-width: 650px) {
    .workspace-stats-grid {
        grid-template-columns: 1fr;
    }
}
}
.ai-summary-card, .winner-card, .feature-impact-card {background:linear-gradient(135deg,rgba(15,35,71,.96),rgba(40,28,92,.92)); border:1px solid rgba(125,211,252,.38); border-radius:24px; padding:26px; margin:18px 0; box-shadow:0 18px 48px rgba(0,0,0,.28); color:#fff;}
.compact-card {padding:22px;}
.ai-summary-title,.winner-title {font-size:24px; font-weight:1000; margin-bottom:18px; color:white;}
.summary-grid {display:grid; grid-template-columns:repeat(4,minmax(0,1fr)); gap:16px; margin:14px 0 18px 0;}
.summary-grid div {background:rgba(15,23,42,.52); border:1px solid rgba(147,197,253,.22); border-radius:16px; padding:16px;}
.summary-grid b {display:block; font-size:28px; color:#FFFFFF; text-shadow:0 0 18px rgba(255,255,255,.18);}
.summary-grid span {display:block; color:#CBD5E1; font-weight:800; margin-top:4px;}
.summary-text,.winner-card p {color:#F8FAFC; font-size:16px; line-height:1.75;}
.premium-list {margin:0; padding-left:20px; color:#F8FAFC; line-height:1.8; font-size:16px;}
.premium-heading {color:#FFFFFF; font-size:26px; font-weight:1000; margin-top:24px;}
.model-score-grid {display:grid; grid-template-columns:repeat(2,minmax(0,1fr)); gap:18px; margin:18px 0;}
.model-score-card {background:linear-gradient(180deg,#1E293B,#0F172A); border:1px solid rgba(147,197,253,.35); border-radius:20px; padding:20px;}
.model-score-title {font-size:22px; font-weight:1000; color:white; margin-bottom:14px;}
.score-line {display:flex; justify-content:space-between; color:#E2E8F0; font-weight:800; margin-top:10px;}
.bar,.feature-impact-bar {height:10px; background:rgba(148,163,184,.25); border-radius:999px; overflow:hidden; margin-top:6px;}
.bar i,.feature-impact-bar i {display:block; height:100%; background:linear-gradient(90deg,#38BDF8,#6366F1,#EC4899); border-radius:999px;}
.feature-impact-row {display:grid; grid-template-columns:240px 1fr 90px; align-items:center; gap:14px; margin:12px 0;}
.feature-impact-name {font-weight:900; color:#FFFFFF;}
.feature-impact-pct {font-weight:900; color:#DBEAFE; text-align:right;}
[data-testid="stDataFrame"] {border-radius:16px; overflow:hidden; border:1px solid rgba(147,197,253,.18);}

</style>
""", unsafe_allow_html=True)


with st.sidebar:
    render_labmind_logo("small")

    st.markdown(
        """
        <hr>
        <div class="sidebar-item">🏠 Dashboard</div>
        <div class="sidebar-item">📊 Overview</div>
        <div class="sidebar-item">🧠 AI Insights</div>
        <div class="sidebar-item">🧹 Auto Clean</div>
        <div class="sidebar-item">📈 Visualizations</div>
        <div class="sidebar-item">🤖 AutoML</div>
        <div class="sidebar-item">🎯 Prediction Playground</div>
        <div class="sidebar-item">💬 Data Chat</div>
        <div class="sidebar-item">📄 Reports</div>
        <hr>
        """,
        unsafe_allow_html=True,
    )

    st.info("Version 3.0 Premium MVP")

def render_labmind_hero():

    """Render the main LabMind hero safely."""
    import base64
    from pathlib import Path

    logo_path = Path(__file__).resolve().parent / "labmind_logo.png"

    if logo_path.exists():
        logo_base64 = base64.b64encode(
            logo_path.read_bytes()
        ).decode("utf-8")

        logo_html = (
            f'<img src="data:image/png;base64,{logo_base64}" '
            f'class="hero-real-logo" alt="LabMind.ai logo">'
        )
    else:
        logo_html = '<span>🧪</span>'

    hero_html = (
        '<section class="labmind-hero">'
        '<div class="hero-brand-inline">'
        '<div class="hero-brand-logo-box">'
        f'{logo_html}'
        '</div>'
        '<div>'
        '<div class="hero-brand-name-new">LabMind.ai</div>'
        '<div class="hero-brand-tagline-new">Premium AI Data Analyst</div>'
        '</div>'
        '</div>'
        '<div class="hero-main-title">'
        'The AI Data Scientist<br>'
        '<span class="hero-title-accent">Built for Every Dataset</span>'
        '</div>'
        '<div class="hero-description">'
        'Upload datasets, clean data automatically, compare machine-learning '
        'models, explain predictions, and generate professional reports—'
        'all from one intelligent workspace.'
        '</div>'
        '<div class="hero-feature-row">'
        '<span>📊 Analyze</span>'
        '<span>🧹 Clean</span>'
        '<span>🤖 Train</span>'
        '<span>🎯 Predict</span>'
        '<span>🧬 Explain</span>'
        '<span>📄 Report</span>'
        '</div>'
        '</section>'
    )

    st.markdown(hero_html, unsafe_allow_html=True)

render_labmind_hero()

dataset_count = 1 if st.session_state.get("df") is not None else 0

model_count = 11

best_accuracy = st.session_state.get("accuracy")
accuracy_text = (
    f"{best_accuracy * 100:.1f}%"
    if isinstance(best_accuracy, (int, float))
    else "—"
)

report_count = 1 if st.session_state.get("report") else 0

workspace_stats_html = (
    '<div class="workspace-stats-grid">'

    '<div class="workspace-stat-card">'
        '<div class="workspace-stat-top">'
            '<div class="workspace-stat-icon">📊</div>'
            '<div class="workspace-stat-label">DATASETS</div>'
        '</div>'
        f'<div class="workspace-stat-value">{dataset_count}</div>'
        '<div class="workspace-stat-caption">Active workspace datasets</div>'
    '</div>'

    '<div class="workspace-stat-card">'
        '<div class="workspace-stat-top">'
            '<div class="workspace-stat-icon">🤖</div>'
            '<div class="workspace-stat-label">MODELS</div>'
        '</div>'
        f'<div class="workspace-stat-value">{model_count}</div>'
        '<div class="workspace-stat-caption">Algorithms available</div>'
    '</div>'

    '<div class="workspace-stat-card">'
        '<div class="workspace-stat-top">'
            '<div class="workspace-stat-icon">⚡</div>'
            '<div class="workspace-stat-label">BEST ACCURACY</div>'
        '</div>'
        f'<div class="workspace-stat-value">{accuracy_text}</div>'
        '<div class="workspace-stat-caption">Current strongest model</div>'
    '</div>'

    '<div class="workspace-stat-card">'
        '<div class="workspace-stat-top">'
            '<div class="workspace-stat-icon">📄</div>'
            '<div class="workspace-stat-label">REPORTS</div>'
        '</div>'
        f'<div class="workspace-stat-value">{report_count}</div>'
        '<div class="workspace-stat-caption">Generated intelligence reports</div>'
    '</div>'

    '</div>'
)

st.markdown(workspace_stats_html, unsafe_allow_html=True)


st.markdown("""
<div class="upload-wrapper">
    <div class="upload-title">Upload Dataset</div>
    <div class="upload-subtitle">
        Drag and drop your CSV file below. LabMind will inspect, score, visualize, clean, model, predict, and report on your data.
    </div>
</div>
""", unsafe_allow_html=True)


uploaded_file = st.file_uploader("📁 Upload or drag your CSV file here", type=["csv"])


if uploaded_file:
    reset_model_state_if_new_file(uploaded_file)

    raw_df = pd.read_csv(uploaded_file)
    df = clean_uploaded_dataframe(raw_df)

    removed_columns = [col for col in raw_df.columns if col not in df.columns]
    st.success("Dataset uploaded successfully!")

    if removed_columns:
        st.info(f"LabMind automatically removed empty/system columns: {removed_columns}")

    total_missing = int(df.isnull().sum().sum())
    total_cells = df.shape[0] * df.shape[1]
    missing_percentage = (total_missing / total_cells) * 100 if total_cells > 0 else 0
    duplicate_rows = int(df.duplicated().sum())
    duplicate_percentage = (duplicate_rows / df.shape[0]) * 100 if df.shape[0] > 0 else 0

    numeric_cols = df.select_dtypes(include=["number"]).columns.tolist()
    text_cols = df.select_dtypes(include=["object"]).columns.tolist()

    completeness_score = max(0, round(100 - missing_percentage))
    uniqueness_score = max(0, round(100 - duplicate_percentage))
    consistency_score = 90 if len(numeric_cols) > 0 else 60
    readiness_score = 90 if df.shape[0] >= 100 and len(numeric_cols) > 0 else 65
    health_score = round((completeness_score + uniqueness_score + consistency_score + readiness_score) / 4)

    recommendations = generate_recommendations(df, health_score, missing_percentage, duplicate_rows)
    cleaning_steps = create_cleaning_steps(df, removed_columns, duplicate_rows, missing_percentage)
    st.session_state["cleaning_steps"] = cleaning_steps

    best_model_display = st.session_state.get("best_model_name", "Not trained")
    best_accuracy_display = f"{st.session_state['best_accuracy'] * 100:.2f}%" if "best_accuracy" in st.session_state else "N/A"
    target_display = st.session_state.get("target_column", "Not selected")

    c1, c2, c3, c4 = st.columns(4)
    for col, title, value in [
        (c1, "Dataset Health", f"{health_score}/100"),
        (c2, "Best Model", best_model_display),
        (c3, "Best Accuracy", best_accuracy_display),
        (c4, "Target Column", target_display),
    ]:
        with col:
            st.markdown(f"""
            <div class="card">
                <div class="card-title">{title}</div>
                <div class="card-value">{value}</div>
            </div>
            """, unsafe_allow_html=True)

    tab1, tab2, tab3, tab4, tab5, tab6, tab7, tab8, tab9, tab10, tab11 = st.tabs([
        "📊 Overview",
        "🧠 AI Insights",
        "🧹 Auto Clean",
        "📈 Visualizations",
        "🤖 AutoML",
        "🎯 Predict",
        "💬 AI Data Chat",
        "🧬 Explainability",
        "📄 Report",
        "🚀 Deploy",
        "💎 Pricing"
    ])

    with tab1:
        st.markdown('<div class="section-title">Dataset Overview</div>', unsafe_allow_html=True)
        render_ai_dataset_summary(df, numeric_cols, text_cols, health_score, missing_percentage, duplicate_rows, df.columns[detect_default_target(df)])

        health_status = (
        "Excellent"
        if health_score >= 85
        else "Good"
        if health_score >= 70
        else "Needs Attention"
        )

        health_status_class = (
        "health-excellent"
        if health_score >= 85
        else "health-good"
        if health_score >= 70
        else "health-warning"
        )

        health_html = textwrap.dedent(f"""    
        <div class="dataset-health-card">
        <div class="dataset-health-header">
            <div>
                <div class="dataset-health-eyebrow">AI DATASET HEALTH</div>
                <div class="dataset-health-title">Executive Dataset Scorecard</div>
            </div>

            <div class="dataset-health-score-wrap">
                <div class="dataset-health-score">{health_score}</div>
                <div class="dataset-health-total">/100</div>
            </div>
        </div>

        <div class="dataset-health-status {health_status_class}">
            {health_status}
        </div>

        <div class="health-metrics-grid">
            <div class="health-metric">
                <div class="health-metric-row">
                    <span>Completeness</span>
                    <strong>{completeness_score}%</strong>
                </div>
                <div class="health-progress">
                    <div style="width:{completeness_score}%"></div>
                </div>
            </div>

            <div class="health-metric">
                <div class="health-metric-row">
                    <span>Uniqueness</span>
                    <strong>{uniqueness_score}%</strong>
                </div>
                <div class="health-progress">
                    <div style="width:{uniqueness_score}%"></div>
                </div>
            </div>

            <div class="health-metric">
                <div class="health-metric-row">
                    <span>Consistency</span>
                    <strong>{consistency_score}%</strong>
                </div>
                <div class="health-progress">
                    <div style="width:{consistency_score}%"></div>
                </div>
            </div>

            <div class="health-metric">
                <div class="health-metric-row">
                    <span>Model Readiness</span>
                    <strong>{readiness_score}%</strong>
                </div>
                <div class="health-progress">
                    <div style="width:{readiness_score}%"></div>
                </div>
            </div>
        </div>
        </div>
        """)

        health_html = "".join(line.strip() for line in health_html.splitlines())

        st.markdown(health_html, unsafe_allow_html=True)
        st.write("### Dataset Preview")
        st.dataframe(df.head(), use_container_width=True)

        st.write("### Dataset Shape")
        st.write(f"Rows: {df.shape[0]}")
        st.write(f"Columns: {df.shape[1]}")

        st.write("### Missing Values")
        st.dataframe(df.isnull().sum(), use_container_width=True)

        st.write("### Column Types")
        st.dataframe(df.dtypes.astype(str), use_container_width=True)

        st.write("### Numeric Statistics")
        if len(numeric_cols) > 0:
            st.dataframe(df[numeric_cols].describe(), use_container_width=True)
        else:
            st.info("No numeric columns found.")

        st.write("### Dataset Comparison")
        comparison_file = st.file_uploader("Optional: upload another CSV to compare against this dataset", type=["csv"], key="comparison_upload")
        if comparison_file is not None:
            comparison_raw = pd.read_csv(comparison_file)
            comparison_df = clean_uploaded_dataframe(comparison_raw)
            comparison_summary = compare_two_datasets(df, comparison_df, uploaded_file.name, comparison_file.name)
            st.dataframe(comparison_summary, use_container_width=True)
            st.download_button(
                "📥 Download Dataset Comparison CSV",
                data=comparison_summary.to_csv(index=False).encode("utf-8"),
                file_name="labmind_dataset_comparison.csv",
                mime="text/csv",
                key="download_dataset_comparison_csv"
            )

    with tab2:
        st.markdown('<div class="section-title">AI Insights & Recommendations</div>', unsafe_allow_html=True)

        a1, a2, a3, a4 = st.columns(4)
        a1.metric("Health Score", f"{health_score}/100")
        a2.metric("Missing Data", f"{missing_percentage:.2f}%")
        a3.metric("Duplicate Rows", duplicate_rows)
        a4.metric("Columns", df.shape[1])

        if health_score >= 85:
            st.success("Excellent dataset quality. Ready for baseline modeling.")
        elif health_score >= 65:
            st.warning("Moderate dataset quality. Cleaning is recommended.")
        else:
            st.error("Low dataset quality. Clean data before modeling.")

        render_premium_recommendations(df, numeric_cols, text_cols, recommendations)

        for rec in recommendations:
            st.markdown(f"""
            <div class="recommendation-box">
                <div class="recommendation-title">🧠 Recommendation</div>
                <div class="recommendation-text">{rec}</div>
            </div>
            """, unsafe_allow_html=True)

    with tab3:
        st.markdown('<div class="section-title">Auto Clean Dataset</div>', unsafe_allow_html=True)

        st.write("LabMind will:")
        st.write("- Remove duplicate rows")
        st.write("- Fill numeric missing values with median")
        st.write("- Fill text missing values with Unknown")

        st.write("### Cleaning Pipeline")
        for step in cleaning_steps:
            st.markdown(f"✅ {step}")

        if st.button("🧹 Auto Clean Dataset"):
            cleaned_df = auto_clean_dataset(df)
            st.session_state["cleaned_df"] = cleaned_df

            after_missing = int(cleaned_df.isnull().sum().sum())
            removed_duplicates = df.shape[0] - cleaned_df.shape[0]

            st.success("Dataset cleaned successfully!")

            x1, x2, x3, x4 = st.columns(4)
            x1.metric("Original Rows", df.shape[0])
            x2.metric("Cleaned Rows", cleaned_df.shape[0])
            x3.metric("Duplicates Removed", removed_duplicates)
            x4.metric("Missing Remaining", after_missing)

            st.write("### Cleaned Dataset Preview")
            st.dataframe(cleaned_df.head(), use_container_width=True)

            csv = cleaned_df.to_csv(index=False).encode("utf-8")
            st.download_button(
                "📥 Download Cleaned CSV",
                data=csv,
                file_name="labmind_cleaned_dataset.csv",
                mime="text/csv",
                key="download_cleaned_csv"
            )

    with tab4:
        st.markdown('<div class="section-title">Visualizations</div>', unsafe_allow_html=True)

        if len(numeric_cols) == 0:
            st.warning("No numeric columns found for visualization.")
        else:
            st.write("### Distribution Chart")
            dist_col = st.selectbox("Select numeric column", numeric_cols, key="dist_col")

            fig, ax = plt.subplots(figsize=(8, 4))
            ax.hist(df[dist_col].dropna(), bins=20)
            ax.set_title(f"Distribution of {dist_col}")
            ax.set_xlabel(dist_col)
            ax.set_ylabel("Frequency")
            st.pyplot(fig)

            if len(numeric_cols) >= 2:
                st.write("### Scatter Plot Builder")
                default_y = "Survived" if "Survived" in numeric_cols else numeric_cols[1]

                x_axis = st.selectbox("X Axis", numeric_cols, key="scatter_x")
                y_axis = st.selectbox(
                    "Y Axis",
                    numeric_cols,
                    index=numeric_cols.index(default_y),
                    key="scatter_y"
                )

                fig2, ax2 = plt.subplots(figsize=(8, 4))
                ax2.scatter(df[x_axis], df[y_axis])
                ax2.set_xlabel(x_axis)
                ax2.set_ylabel(y_axis)
                ax2.set_title(f"{x_axis} vs {y_axis}")
                st.pyplot(fig2)

            st.write("### Correlation Heatmap")
            if len(numeric_cols) >= 2:
                corr = df[numeric_cols].corr()
                fig3, ax3 = plt.subplots(figsize=(10, 6))
                im = ax3.imshow(corr)
                ax3.set_xticks(range(len(corr.columns)))
                ax3.set_yticks(range(len(corr.columns)))
                ax3.set_xticklabels(corr.columns, rotation=45, ha="right")
                ax3.set_yticklabels(corr.columns)
                ax3.set_title("Correlation Heatmap")
                fig3.colorbar(im)
                st.pyplot(fig3)

            st.write("### Outlier Detection")
            outlier_col = st.selectbox("Choose column for boxplot", numeric_cols, key="outlier_col")
            q1 = df[outlier_col].quantile(0.25)
            q3 = df[outlier_col].quantile(0.75)
            iqr = q3 - q1
            lower = q1 - 1.5 * iqr
            upper = q3 + 1.5 * iqr
            outliers = df[(df[outlier_col] < lower) | (df[outlier_col] > upper)]

            fig4, ax4 = plt.subplots(figsize=(8, 4))
            ax4.boxplot(df[outlier_col].dropna(), vert=False)
            ax4.set_title(f"Boxplot for {outlier_col}")
            st.pyplot(fig4)

            st.info(f"Potential outliers in {outlier_col}: {outliers.shape[0]}")

    with tab5:
        st.markdown('<div class="section-title">AutoML Model Training</div>', unsafe_allow_html=True)

        default_target = detect_default_target(df)

        target_column = st.selectbox(
            "Choose target column",
            df.columns,
            index=default_target,
            key="automl_target"
        )

        run_hyperparameter_tuning = st.checkbox("🔧 Run lightweight hyperparameter optimization after baseline training", value=False)

        if st.button("🚀 Train & Compare Models"):
            if target_column.lower() in ["passengerid", "id"]:
                st.error("Do not use an ID column as the target. Choose something meaningful like Survived.")
            else:
                data, X, y, category_maps, numeric_defaults, target_label_map = prepare_ml_data(df, target_column)

                if y.nunique() < 2:
                    st.error(f"This target column has only one class after preprocessing: {list(y.unique())}. Choose a different target column.")
                    st.stop()

                stratify_y = None
                if y.nunique() <= 20 and y.value_counts().min() >= 2:
                    stratify_y = y

                X_train, X_test, y_train, y_test = train_test_split(
                    X,
                    y,
                    test_size=0.2,
                    random_state=42,
                    stratify=stratify_y
                )

                models = build_classification_models(y)

                st.info(f"LabMind will compare {len(models)} classification models for this target.")

                results = []
                trained_models = {}

                for name, model in models.items():
                    try:
                        model.fit(X_train, y_train)
                        preds = model.predict(X_test)
                        acc = accuracy_score(y_test, preds)
                        f1 = f1_score(y_test, preds, average="weighted", zero_division=0)

                        cv_mean = np.nan
                        cv_std = np.nan
                        try:
                            cv_splits = min(5, int(y.value_counts().min())) if y.nunique() <= 20 else 3
                            if cv_splits >= 2:
                                cv = StratifiedKFold(n_splits=cv_splits, shuffle=True, random_state=42)
                                cv_scores = cross_val_score(model, X, y, cv=cv, scoring="accuracy")
                                cv_mean = float(cv_scores.mean())
                                cv_std = float(cv_scores.std())
                        except Exception:
                            pass

                        results.append({"Model": name, "Accuracy": acc, "Weighted F1": f1, "CV Mean": cv_mean, "CV Std": cv_std})
                        trained_models[name] = {"model": model, "predictions": preds, "accuracy": acc, "weighted_f1": f1, "cv_mean": cv_mean, "cv_std": cv_std}
                    except Exception as e:
                        st.warning(f"{name} skipped: {e}")

                if run_hyperparameter_tuning and len(results) > 0:
                    with st.spinner("Optimizing top baseline models..."):
                        results, trained_models, tuning_log = tune_top_models(
                            results, trained_models, X_train, y_train, X_test, y_test, top_k=3
                        )
                        st.session_state["tuning_results"] = tuning_log
                    if tuning_log:
                        st.write("### 🔧 Hyperparameter Optimization Results")
                        st.dataframe(pd.DataFrame(tuning_log), use_container_width=True)

                if len(results) == 0:
                    st.error("No models could be trained. Try a different target column.")
                else:
                    results_df = pd.DataFrame(results).sort_values(by="Accuracy", ascending=False)
                    best_model_name = results_df.iloc[0]["Model"]
                    best_accuracy = results_df.iloc[0]["Accuracy"]
                    best_model = trained_models[best_model_name]["model"]
                    best_predictions = trained_models[best_model_name]["predictions"]

                    st.session_state["best_model"] = best_model
                    st.session_state["best_model_name"] = best_model_name
                    st.session_state["best_accuracy"] = best_accuracy
                    st.session_state["target_column"] = target_column
                    st.session_state["feature_columns"] = list(X.columns)
                    st.session_state["category_maps"] = category_maps
                    st.session_state["numeric_defaults"] = numeric_defaults
                    st.session_state["target_label_map"] = target_label_map
                    st.session_state["trained_models"] = trained_models
                    st.session_state["leaderboard"] = results_df.copy()
                    st.session_state["X_train"] = X_train
                    st.session_state["X_test"] = X_test
                    st.session_state["y_train"] = y_train
                    st.session_state["y_test"] = y_test

                    st.success("Models trained successfully!")

                    m1, m2 = st.columns(2)
                    m1.metric("Best Model", best_model_name)
                    m2.metric("Best Accuracy", f"{best_accuracy * 100:.2f}%")

                    st.write("### 🏆 Model Leaderboard")
                    leaderboard = results_df.copy().sort_values(by="Accuracy", ascending=False).reset_index(drop=True)
                    leaderboard_display = leaderboard.copy()
                    leaderboard_display["Accuracy"] = leaderboard_display["Accuracy"].apply(lambda x: f"{x * 100:.2f}%")
                    if "Weighted F1" in leaderboard_display.columns:
                        leaderboard_display["Weighted F1"] = leaderboard_display["Weighted F1"].apply(lambda x: f"{x * 100:.2f}%")
                    if "CV Mean" in leaderboard_display.columns:
                        leaderboard_display["CV Mean"] = leaderboard_display["CV Mean"].apply(lambda x: "N/A" if pd.isna(x) else f"{x * 100:.2f}%")
                    if "CV Std" in leaderboard_display.columns:
                        leaderboard_display["CV Std"] = leaderboard_display["CV Std"].apply(lambda x: "N/A" if pd.isna(x) else f"±{x * 100:.2f}%")
                    leaderboard_display.insert(0, "Rank", (["🥇", "🥈", "🥉"] + [f"#{i}" for i in range(4, len(leaderboard_display) + 1)])[:len(leaderboard_display)])
                    st.dataframe(leaderboard_display, use_container_width=True)

                    runner_up_text = ""
                    if len(leaderboard) > 1:
                        runner_up_text = f" The runner-up was {leaderboard.iloc[1]['Model']} with {leaderboard.iloc[1]['Accuracy'] * 100:.2f}% accuracy."

                    st.markdown(f"""
                    <div class="model-highlight">
                        🏆 <b>{best_model_name}</b> is currently the strongest model with <b>{best_accuracy * 100:.2f}%</b> accuracy.{runner_up_text}<br><br>
                        🧠 LabMind recommendation: use this as your baseline model, then validate it on a separate dataset before deployment.
                    </div>
                    """, unsafe_allow_html=True)

                    render_ai_model_card(best_model_name, best_accuracy, leaderboard)
                    render_model_comparison_dashboard(leaderboard)
                    render_model_comparison_charts(leaderboard)

                    cm = confusion_matrix(y_test, best_predictions)
                    st.session_state["confusion_matrix"] = cm

                    st.write("### Confusion Matrix")
                    render_confusion_summary(cm)

                    fig_cm, ax_cm = plt.subplots(figsize=(6, 4))
                    im_cm = ax_cm.imshow(cm)
                    ax_cm.set_title("Confusion Matrix")
                    ax_cm.set_xlabel("Predicted")
                    ax_cm.set_ylabel("Actual")
                    ax_cm.set_xticks(range(cm.shape[1]))
                    ax_cm.set_yticks(range(cm.shape[0]))
                    if target_label_map:
                        labels = [str(target_label_map.get(i, i)) for i in range(cm.shape[0])]
                        ax_cm.set_xticklabels(labels)
                        ax_cm.set_yticklabels(labels)

                    for i in range(cm.shape[0]):
                        for j in range(cm.shape[1]):
                            ax_cm.text(j, i, cm[i, j], ha="center", va="center", fontsize=14, fontweight="bold")

                    fig_cm.colorbar(im_cm, ax=ax_cm)
                    st.pyplot(fig_cm)

                    if hasattr(best_model, "predict_proba") and len(pd.Series(y_test).unique()) == 2:
                        st.write("### ROC Curve")
                        y_prob = best_model.predict_proba(X_test)[:, 1]
                        fpr, tpr, _ = roc_curve(y_test, y_prob)
                        roc_auc = auc(fpr, tpr)
                        st.session_state["roc_auc"] = roc_auc
                        render_auc_summary(roc_auc)

                        fig_roc, ax_roc = plt.subplots(figsize=(7, 5))
                        ax_roc.plot(fpr, tpr, label=f"AUC = {roc_auc:.3f}")
                        ax_roc.plot([0, 1], [0, 1], linestyle="--")
                        ax_roc.set_xlabel("False Positive Rate")
                        ax_roc.set_ylabel("True Positive Rate")
                        ax_roc.set_title("ROC Curve")
                        ax_roc.legend(loc="lower right")
                        st.pyplot(fig_roc)

                    st.write("### ROC Curves Across Models")
                    render_all_model_roc_curves(trained_models, X_test, y_test)

                    st.write("### Classification Report")
                    report_df = make_readable_classification_report(y_test, best_predictions, target_label_map)
                    st.dataframe(report_df, use_container_width=True)

                    weighted_precision = report_df.loc["weighted avg", "precision"] if "weighted avg" in report_df.index else None
                    weighted_recall = report_df.loc["weighted avg", "recall"] if "weighted avg" in report_df.index else None
                    weighted_f1 = report_df.loc["weighted avg", "f1-score"] if "weighted avg" in report_df.index else None

                    st.write("### Model Performance Summary")
                    perf1, perf2, perf3, perf4 = st.columns(4)
                    perf1.metric("Accuracy", f"{best_accuracy * 100:.2f}%")
                    perf2.metric("Weighted Precision", f"{weighted_precision * 100:.2f}%" if weighted_precision is not None else "N/A")
                    perf3.metric("Weighted Recall", f"{weighted_recall * 100:.2f}%" if weighted_recall is not None else "N/A")
                    perf4.metric("Weighted F1", f"{weighted_f1 * 100:.2f}%" if weighted_f1 is not None else "N/A")

                    best_row = leaderboard[leaderboard["Model"] == best_model_name].iloc[0]
                    st.write("### Cross-Validation Stability")
                    cv1, cv2 = st.columns(2)
                    cv_mean_value = best_row.get("CV Mean", np.nan)
                    cv_std_value = best_row.get("CV Std", np.nan)
                    cv1.metric("CV Mean Accuracy", "N/A" if pd.isna(cv_mean_value) else f"{cv_mean_value * 100:.2f}%")
                    cv2.metric("CV Variation", "N/A" if pd.isna(cv_std_value) else f"±{cv_std_value * 100:.2f}%")
                    if not pd.isna(cv_mean_value):
                        st.markdown(f"""
                        <div class="model-highlight">
                            ✅ Cross-validation checks whether the model stays reliable across multiple data splits. 
                            The best model averaged <b>{cv_mean_value * 100:.2f}%</b> accuracy with variation of <b>±{cv_std_value * 100:.2f}%</b>.
                        </div>
                        """, unsafe_allow_html=True)

                    st.session_state["classification_report_df"] = report_df.copy()

                    st.download_button(
                        "📊 Download Metrics CSV",
                        data=report_df.to_csv().encode("utf-8"),
                        file_name="labmind_model_metrics.csv",
                        mime="text/csv",
                        key="download_metrics_csv_automl"
                    )

                    st.write("### Feature Importance")
                    importance_values = get_model_importance(best_model, X.columns)

                    importance = pd.DataFrame({
                        "Feature": X.columns,
                        "Importance": importance_values
                    }).sort_values(by="Importance", ascending=False)

                    importance_top = prepare_top_importance_table(importance, top_n=10)
                    st.dataframe(importance_top, use_container_width=True)

                    fig5, ax5 = plt.subplots(figsize=(10, 5))
                    top_plot = importance.head(15).sort_values("Importance", ascending=True)
                    ax5.barh(top_plot["Feature"], top_plot["Importance"])
                    ax5.set_xlabel("Importance")
                    ax5.set_ylabel("Feature")
                    ax5.set_title("Top 15 Feature Importance")
                    st.pyplot(fig5)

                    top_feature = importance.iloc[0]["Feature"]
                    top_value = importance.iloc[0]["Importance"]

                    report_text = f"""
LabMind.ai Analysis Report

Executive Summary:
- Dataset contains {df.shape[0]} rows and {df.shape[1]} columns.
- Dataset Health Score: {health_score}/100
- Completeness Score: {completeness_score}/100
- Uniqueness Score: {uniqueness_score}/100
- Consistency Score: {consistency_score}/100
- Model Readiness Score: {readiness_score}/100

Dataset Quality:
- Missing Values: {total_missing}
- Missing Data Percentage: {missing_percentage:.2f}%
- Duplicate Rows: {duplicate_rows}
- Numeric Columns: {len(numeric_cols)}
- Text Columns: {len(text_cols)}

Machine Learning Summary:
- Target Column: {target_column}
- Best Model: {best_model_name}
- Best Accuracy: {best_accuracy * 100:.2f}%

Top Feature:
- {top_feature}
- Importance Score: {top_value:.4f}

Recommendations:
{chr(10).join("- " + r for r in recommendations)}
"""
                    st.session_state["report"] = report_text

                    executive_summary = generate_executive_summary(
                        df=df,
                        health_score=health_score,
                        missing_percentage=missing_percentage,
                        duplicate_rows=duplicate_rows,
                        best_model_name=best_model_name,
                        best_accuracy=best_accuracy,
                        target_column=target_column,
                        top_feature=top_feature
                    )
                    st.session_state["executive_summary"] = executive_summary

                    st.markdown("### 🧠 AI Executive Summary")
                    summary_html = executive_summary.replace("\n", "<br>")
                    st.markdown(f"""
                    <div class="insight-card">
                        <h3>🧠 AI Executive Summary</h3>
                        <p>{summary_html}</p>
                    </div>
                    """, unsafe_allow_html=True)

                    st.write("### 🔍 Model Explainability")
                    explainability_df = importance.head(8).copy()
                    explainability_df["Explanation"] = explainability_df["Feature"].apply(
                        lambda x: f"{x} contributed to the model based on learned feature importance."
                    )
                    st.dataframe(explainability_df, use_container_width=True)

                    fig_exp, ax_exp = plt.subplots(figsize=(10, 5))
                    ax_exp.barh(explainability_df["Feature"], explainability_df["Importance"])
                    ax_exp.set_title("Top Feature Contributions")
                    ax_exp.set_xlabel("Importance")
                    ax_exp.invert_yaxis()
                    st.pyplot(fig_exp)

                    st.write("### 🧬 Universal Model Explainability")
                    explain_df, explain_method, explain_note = compute_universal_explainability(
                        best_model, X_train, X_test, y_test, X.columns
                    )

                    ai_feature_summary = build_feature_ai_summary(explain_df, explain_method, best_model_name, target_column)
                    st.session_state["model_explainability_summary"] = ai_feature_summary
                    st.session_state["explain_df"] = explain_df.copy()
                    st.session_state["explain_method"] = explain_method

                    st.markdown(f"""
                    <div class="model-highlight">
                        🧠 <b>{explain_method}</b><br>
                        {explain_note}<br><br>
                        <b>AI Summary:</b> {ai_feature_summary}
                    </div>
                    """, unsafe_allow_html=True)

                    explain_top = prepare_top_importance_table(explain_df, top_n=10)
                    st.dataframe(explain_top, use_container_width=True)

                    fig_uni, ax_uni = plt.subplots(figsize=(10, 5))
                    top_explain = explain_df.head(15).sort_values("Importance", ascending=True)
                    ax_uni.barh(top_explain["Feature"], top_explain["Importance"])
                    ax_uni.set_xlabel("Importance")
                    ax_uni.set_title(f"{explain_method} - Top 15 Feature Impact")
                    st.pyplot(fig_uni)

                    st.download_button(
                        "📥 Download Explainability CSV",
                        data=explain_df.to_csv(index=False).encode("utf-8"),
                        file_name="labmind_explainability.csv",
                        mime="text/csv",
                        key="download_explainability_csv_automl"
                    )

                    model_buffer = BytesIO()
                    joblib.dump(best_model, model_buffer)
                    model_buffer.seek(0)

                    st.download_button(
                        "📦 Download Best Model (.pkl)",
                        data=model_buffer,
                        file_name="labmind_best_model.pkl",
                        mime="application/octet-stream",
                        key="download_best_model_automl"
                    )

    with tab6:
        st.markdown('<div class="section-title">Prediction Playground</div>', unsafe_allow_html=True)

        if "best_model" not in st.session_state:
            st.info("Train models first in the AutoML tab to unlock prediction.")
        else:
            st.success(
                f"Using best model: {st.session_state['best_model_name']} "
                f"({st.session_state['best_accuracy'] * 100:.2f}% accuracy)"
            )

            model = st.session_state["best_model"]
            feature_columns = st.session_state["feature_columns"]
            category_maps = st.session_state["category_maps"]
            numeric_defaults = st.session_state["numeric_defaults"]
            target_label_map = st.session_state.get("target_label_map")

            input_data = {}

            for feature in feature_columns:
                if feature in category_maps:
                    options = category_maps[feature]
                    selected = st.selectbox(feature, options, key=f"predict_{feature}")
                    input_data[feature] = options.index(selected)
                else:
                    default_value = float(numeric_defaults.get(feature, 0))
                    value = st.number_input(feature, value=default_value, key=f"predict_{feature}")
                    input_data[feature] = value

            if st.button("🎯 Predict"):
                input_df = pd.DataFrame([input_data])
                input_df = input_df[feature_columns]

                prediction = model.predict(input_df)[0]
                decoded_prediction = target_label_map.get(prediction, prediction) if target_label_map else prediction

                p1, p2 = st.columns(2)
                p1.metric("Predicted Target", decoded_prediction)

                if hasattr(model, "predict_proba"):
                    probs = model.predict_proba(input_df)[0]
                    confidence = max(probs) * 100
                    p2.metric("Confidence", f"{confidence:.2f}%")
                    render_probability_bars(probs, model.classes_, target_label_map)
                else:
                    p2.metric("Confidence", "N/A")

                confidence_text = "N/A"
                if hasattr(model, "predict_proba"):
                    confidence_text = f"{confidence:.2f}%"

                history_row = {
                    "Model": st.session_state["best_model_name"],
                    "Target": st.session_state["target_column"],
                    "Prediction": decoded_prediction,
                    "Confidence": confidence_text
                }
                st.session_state.setdefault("prediction_history", []).append(history_row)

                st.success(f"LabMind predicts {st.session_state['target_column']} = {decoded_prediction}")
                render_prediction_reason_and_similar(input_df, decoded_prediction, confidence_text)

            if st.session_state.get("prediction_history"):
                st.write("### Prediction History")
                history_df = pd.DataFrame(st.session_state["prediction_history"])
                st.dataframe(history_df, use_container_width=True)
                st.download_button(
                    "📥 Download Prediction History CSV",
                    data=history_df.to_csv(index=False).encode("utf-8"),
                    file_name="labmind_prediction_history.csv",
                    mime="text/csv",
                    key="download_prediction_history_csv"
                )

    with tab7:
        st.markdown('<div class="section-title">Chat With Your Dataset</div>', unsafe_allow_html=True)

        st.markdown("""
        <div class="faq-grid">
            <div class="faq-card">How many rows are there?</div>
            <div class="faq-card">Which columns have missing values?</div>
            <div class="faq-card">What is the dataset health score?</div>
            <div class="faq-card">Are there duplicate rows?</div>
            <div class="faq-card">What are the numeric columns?</div>
            <div class="faq-card">What is the best model?</div>
            <div class="faq-card">What is the best accuracy?</div>
            <div class="faq-card">Summarize this dataset.</div>
        </div>
        """, unsafe_allow_html=True)

        question = st.text_input("Ask LabMind a question about your dataset")

        if question:
            answer = answer_dataset_question(
                question,
                df,
                numeric_cols,
                text_cols,
                health_score,
                total_missing,
                missing_percentage,
                duplicate_rows,
                recommendations
            )

            st.markdown(f"""
            <div class="chat-box">
                <div class="recommendation-title">💬 LabMind Answer</div>
                <div class="chat-answer">{answer}</div>
            </div>
            """, unsafe_allow_html=True)

    with tab8:
        st.markdown('<div class="section-title">Explainability Dashboard</div>', unsafe_allow_html=True)

        if "explain_df" not in st.session_state:
            st.info("Train a model first to unlock the explainability dashboard.")
        else:
            explain_df = st.session_state["explain_df"]
            explain_method = st.session_state.get("explain_method", "Model Explainability")
            st.markdown(f"""
            <div class="model-highlight">
                🧬 <b>{explain_method}</b><br>
                {st.session_state.get("model_explainability_summary", "LabMind identified the strongest model drivers.")}
            </div>
            """, unsafe_allow_html=True)

            st.write("### Top 10 Model Drivers")
            st.dataframe(prepare_top_importance_table(explain_df, 10), use_container_width=True)
            render_feature_impact_blocks(explain_df, top_n=8)

            st.write("### Top 15 Feature Impact Chart")
            fig_dash, ax_dash = plt.subplots(figsize=(10, 5))
            top_dash = explain_df.head(15).sort_values("Importance", ascending=True)
            ax_dash.barh(top_dash["Feature"], top_dash["Importance"])
            ax_dash.set_title(f"{explain_method} - Top 15 Feature Impact")
            ax_dash.set_xlabel("Importance")
            st.pyplot(fig_dash)

            if "best_model" in st.session_state and "X_train" in st.session_state and "X_test" in st.session_state:
                st.write("### SHAP Explainability")
                render_shap_dashboard(
                    st.session_state["best_model"],
                    st.session_state["X_train"],
                    st.session_state["X_test"],
                    st.session_state.get("feature_columns", list(explain_df["Feature"]))
                )

            if "confusion_matrix" in st.session_state:
                st.write("### Confusion Matrix Summary")
                render_confusion_summary(st.session_state["confusion_matrix"])
                cm_dash = st.session_state["confusion_matrix"]
                fig_cm_dash, ax_cm_dash = plt.subplots(figsize=(6, 4))
                im_dash = ax_cm_dash.imshow(cm_dash)
                ax_cm_dash.set_title("Confusion Matrix Heatmap")
                ax_cm_dash.set_xlabel("Predicted")
                ax_cm_dash.set_ylabel("Actual")
                for i in range(cm_dash.shape[0]):
                    for j in range(cm_dash.shape[1]):
                        ax_cm_dash.text(j, i, cm_dash[i, j], ha="center", va="center", fontsize=14, fontweight="bold")
                fig_cm_dash.colorbar(im_dash, ax=ax_cm_dash)
                st.pyplot(fig_cm_dash)

            if "roc_auc" in st.session_state:
                st.write("### ROC / AUC Summary")
                render_auc_summary(st.session_state["roc_auc"])

            st.download_button(
                "📥 Download Explainability Dashboard CSV",
                data=explain_df.to_csv(index=False).encode("utf-8"),
                file_name="labmind_explainability_dashboard.csv",
                mime="text/csv",
                key="download_explainability_dashboard_csv"
            )

    with tab10:
        st.markdown('<div class="section-title">Deployment Center</div>', unsafe_allow_html=True)

        if "best_model" not in st.session_state:
            st.info("Train models first to unlock deployment options.")
        else:
            st.markdown(f"""
            <div class="model-highlight">
                🚀 <b>Deployment-ready baseline:</b> {st.session_state.get("best_model_name")}<br>
                Accuracy: <b>{st.session_state.get("best_accuracy", 0) * 100:.2f}%</b><br>
                Target: <b>{st.session_state.get("target_column")}</b>
            </div>
            """, unsafe_allow_html=True)

            st.write("### Model Repository")
            if "leaderboard" in st.session_state:
                repo_df = st.session_state["leaderboard"].copy()
                display_repo = repo_df.copy()
                for col in ["Accuracy", "Weighted F1", "CV Mean", "CV Std"]:
                    if col in display_repo.columns:
                        display_repo[col] = display_repo[col].apply(lambda x: "N/A" if pd.isna(x) else f"{x * 100:.2f}%")
                st.dataframe(display_repo, use_container_width=True)
                st.write("### Model Cards")
                render_model_cards(st.session_state.get("trained_models"), repo_df)

            single_buffer = BytesIO()
            joblib.dump(st.session_state["best_model"], single_buffer)
            single_buffer.seek(0)
            st.download_button(
                "📦 Download Best Model (.pkl)",
                data=single_buffer,
                file_name="labmind_best_model.pkl",
                mime="application/octet-stream",
                key="download_best_model_deploy"
            )

            if "trained_models" in st.session_state:
                zip_buffer = create_model_zip(st.session_state["trained_models"])
                st.download_button(
                    "🗂️ Download All Trained Models (.zip)",
                    data=zip_buffer,
                    file_name="labmind_model_repository.zip",
                    mime="application/zip",
                    key="download_model_repository_zip"
                )

            deploy_metadata = {
                "best_model": st.session_state.get("best_model_name"),
                "target_column": st.session_state.get("target_column"),
                "accuracy": st.session_state.get("best_accuracy"),
                "feature_columns": st.session_state.get("feature_columns"),
                "target_label_map": st.session_state.get("target_label_map"),
                "numeric_defaults": st.session_state.get("numeric_defaults"),
                "category_maps": st.session_state.get("category_maps"),
            }
            st.download_button(
                "🧾 Download Deployment Metadata",
                data=pd.Series(deploy_metadata).to_json(indent=2).encode("utf-8"),
                file_name="labmind_deployment_metadata.json",
                mime="application/json",
                key="download_deployment_metadata_json"
            )

            st.write("### One-click Deployment Targets")
            st.markdown("""
            <div class="ai-summary-card compact-card">
                <ul class="premium-list">
                    <li>🚀 Streamlit Cloud for a quick public demo</li>
                    <li>☁️ Render/Railway for app hosting</li>
                    <li>🧠 Hugging Face Spaces for ML portfolio demos</li>
                    <li>🏢 AWS/Azure/GCP after validation</li>
                </ul>
            </div>
            """, unsafe_allow_html=True)

            st.write("### Deployment Checklist")
            for item in [
                "Validate performance with cross-validation",
                "Test prediction playground with realistic examples",
                "Download model and metadata together",
                "Save the exact training dataset version",
                "Avoid deploying before external validation",
            ]:
                st.markdown(f"✅ {item}")

    with tab9:
        st.markdown('<div class="section-title">Executive Report</div>', unsafe_allow_html=True)

        if "report" in st.session_state:
            st.text_area("Generated Report", st.session_state["report"], height=420)
            st.markdown("""
            <div class="winner-card">
                <div class="winner-title">📄 15-Section Executive Report Includes</div>
                <ul class="premium-list">
                    <li>Cover + executive summary</li><li>Dataset quality scorecard</li><li>Cleaning pipeline</li>
                    <li>EDA + visualization findings</li><li>Model leaderboard</li><li>Winner model card</li>
                    <li>Cross-validation stability</li><li>Confusion matrix summary</li><li>ROC/AUC summary when available</li>
                    <li>Explainability drivers</li><li>Prediction history</li><li>Deployment checklist</li>
                </ul>
            </div>
            """, unsafe_allow_html=True)

            st.download_button(
                "Download TXT Report",
                data=st.session_state["report"],
                file_name="labmind_report.txt",
                mime="text/plain",
                key="download_txt_report"
            )

            pdf_file = create_pdf_report(st.session_state["report"])
            st.download_button(
                "Download PDF Report",
                data=pdf_file,
                file_name="labmind_report.pdf",
                mime="application/pdf",
                key="download_pdf_report_basic"
            )

            pro_pdf = create_professional_pdf_report(
                st.session_state["report"],
                leaderboard_df=st.session_state.get("leaderboard"),
                explain_df=st.session_state.get("explain_df"),
                prediction_history=st.session_state.get("prediction_history")
            )
            st.download_button(
                "Download Professional Executive PDF",
                data=pro_pdf,
                file_name="labmind_professional_executive_report.pdf",
                mime="application/pdf",
                key="download_pdf_report_professional"
            )

            ppt_file = create_ppt_report(st.session_state["report"])
            st.download_button(
                "Download PowerPoint Report",
                data=ppt_file,
                file_name="labmind_report.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                key="download_ppt_report"
            )
        else:
            st.info("Train models first to generate a report.")

        with tab11:
             render_pricing_page()

else:
    st.info("Upload a CSV file to begin.")